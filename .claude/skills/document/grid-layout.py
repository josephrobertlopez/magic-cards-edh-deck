#!/usr/bin/env python3
"""
Document Grid Layout Skill Executor

Generic grid layout generation for any set of images.
Works with card proxies, product catalogs, photo albums, employee directories, etc.
"""

import sys
import json
import argparse
from typing import Dict, Any, List
from pathlib import Path
import math
from pptx import Presentation
from pptx.util import Inches


def create_grid_layout(
    image_paths: List[str],
    rows: int,
    columns: int,
    output_path: str,
    slide_width: float = 10.0,
    slide_height: float = 7.5,
    margin: float = 0.5,
    spacing: float = 0.1,
    placeholder_text: str = "Image Not Available"
) -> Dict[str, Any]:
    """
    Create presentation with images arranged in grid layout.

    Args:
        image_paths: List of image file paths
        rows: Number of rows per slide
        columns: Number of columns per slide
        output_path: Path to save PPTX file
        slide_width: Slide width in inches
        slide_height: Slide height in inches
        margin: Margin around grid in inches
        spacing: Spacing between cells in inches
        placeholder_text: Text for missing images

    Returns:
        Dictionary with keys:
        - path: Full path to generated PPTX
        - slide_count: Number of slides created
        - images_per_slide: Max images per slide (rows × columns)
        - missing_images_count: Count of missing/invalid images

    Raises:
        Exception: On validation or generation failure
    """
    # Validate inputs
    if rows <= 0 or columns <= 0:
        raise Exception("INVALID_GRID: rows and columns must be > 0")

    if not image_paths:
        raise Exception("EMPTY_INPUT: image_paths cannot be empty")

    # Calculate layout
    images_per_slide = rows * columns
    slide_count = math.ceil(len(image_paths) / images_per_slide)

    # Calculate cell dimensions
    available_width = slide_width - (2 * margin) - ((columns - 1) * spacing)
    available_height = slide_height - (2 * margin) - ((rows - 1) * spacing)
    cell_width = available_width / columns
    cell_height = available_height / rows

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)

    missing_count = 0

    # Process images in batches (one batch per slide)
    for slide_num in range(slide_count):
        # Use blank layout (index 6 is typically blank)
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_layout)

        # Get batch of images for this slide
        start_idx = slide_num * images_per_slide
        end_idx = min(start_idx + images_per_slide, len(image_paths))
        batch = image_paths[start_idx:end_idx]

        # Place images in grid
        for idx, img_path in enumerate(batch):
            row = idx // columns
            col = idx % columns

            # Calculate position
            left = margin + (col * (cell_width + spacing))
            top = margin + (row * (cell_height + spacing))

            # Add image or placeholder
            img_file = Path(img_path)
            if img_file.exists():
                try:
                    slide.shapes.add_picture(
                        str(img_file),
                        Inches(left),
                        Inches(top),
                        width=Inches(cell_width),
                        height=Inches(cell_height)
                    )
                except Exception:
                    # Invalid image format
                    add_text_placeholder(slide, left, top, cell_width, cell_height, placeholder_text)
                    missing_count += 1
            else:
                # Missing file
                add_text_placeholder(slide, left, top, cell_width, cell_height, placeholder_text)
                missing_count += 1

    # Create output directory if needed
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)

    # Save presentation
    prs.save(str(output_file))

    return {
        "path": str(output_file),
        "slide_count": slide_count,
        "images_per_slide": images_per_slide,
        "missing_images_count": missing_count
    }


def add_text_placeholder(slide, left: float, top: float, width: float, height: float, text: str):
    """Add centered text box as placeholder for missing image."""
    textbox = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    # Center align
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = 1  # PP_ALIGN.CENTER


def main():
    """CLI entry point for skill executor."""
    parser = argparse.ArgumentParser(description="Create grid layout presentation from images")
    parser.add_argument("--image-paths", type=json.loads, required=True,
                       help="JSON list of image paths")
    parser.add_argument("--rows", type=int, required=True,
                       help="Number of rows per slide")
    parser.add_argument("--columns", type=int, required=True,
                       help="Number of columns per slide")
    parser.add_argument("--output-path", required=True,
                       help="Output PPTX file path")
    parser.add_argument("--slide-width", type=float, default=10.0,
                       help="Slide width in inches")
    parser.add_argument("--slide-height", type=float, default=7.5,
                       help="Slide height in inches")
    parser.add_argument("--margin", type=float, default=0.5,
                       help="Margin in inches")
    parser.add_argument("--spacing", type=float, default=0.1,
                       help="Cell spacing in inches")
    parser.add_argument("--placeholder-text", default="Image Not Available",
                       help="Text for missing images")

    args = parser.parse_args()

    try:
        result = create_grid_layout(
            image_paths=args.image_paths,
            rows=args.rows,
            columns=args.columns,
            output_path=args.output_path,
            slide_width=args.slide_width,
            slide_height=args.slide_height,
            margin=args.margin,
            spacing=args.spacing,
            placeholder_text=args.placeholder_text
        )

        print(json.dumps(result))
        sys.exit(0)

    except Exception as e:
        error_result = {
            "success": False,
            "error": str(e),
            "path": None,
            "slide_count": 0
        }
        print(json.dumps(error_result), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
