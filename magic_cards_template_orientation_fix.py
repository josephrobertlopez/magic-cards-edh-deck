#!/usr/bin/env python3
"""
Magic Cards Template Automation - ORIENTATION AWARE VERSION
Respects the original template's mixed orientation design:
- 2 vertical cards (left side)  
- 6 horizontal cards (right side, 3x2 grid)

Based on visual analysis of template PNG.
"""

import os
import glob
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

def analyze_card_slot_orientation(shape, slide_num):
    """
    Determine if a card slot should be vertical or horizontal
    based on its position and the template design
    """
    width_inches = shape.width / 914400
    height_inches = shape.height / 914400
    left_inches = shape.left / 914400
    top_inches = shape.top / 914400
    
    print(f"    📐 Slot analysis: {width_inches:.2f}×{height_inches:.2f} at ({left_inches:.2f},{top_inches:.2f})")
    
    # Based on visual analysis of the template PNG:
    # Slide 1: 2 vertical slots on the left side  
    # Slide 2: 6 horizontal slots in 3x2 grid on the right
    
    if slide_num == 1:
        # First slide should have vertical orientations
        return 'vertical'
    elif slide_num == 2:
        # Second slide should have horizontal orientations  
        return 'horizontal'
    else:
        # Default for any additional slides
        return 'vertical'

def resize_card_for_slot(card_path, slot_width_inches, slot_height_inches, orientation):
    """
    Resize a Magic card image to fit a slot while maintaining aspect ratio
    and respecting the intended orientation
    """
    try:
        with Image.open(card_path) as img:
            # Standard Magic card aspect ratio is 2.5:3.5 (width:height) = 0.714
            card_aspect = 2.5 / 3.5
            
            if orientation == 'horizontal':
                # For horizontal slots, we want the card rotated 90 degrees
                # So we flip the aspect ratio: height becomes width
                target_aspect = 3.5 / 2.5  # 1.4 (horizontal orientation)
                
                # Rotate the image 90 degrees clockwise
                img = img.rotate(-90, expand=True)
                
                # Calculate size to fit the horizontal slot
                slot_aspect = slot_width_inches / slot_height_inches
                
                if target_aspect > slot_aspect:
                    # Image is wider than slot - fit to width
                    new_width = slot_width_inches
                    new_height = slot_width_inches / target_aspect
                else:
                    # Image is taller than slot - fit to height  
                    new_height = slot_height_inches
                    new_width = slot_height_inches * target_aspect
                    
            else:  # vertical orientation
                # For vertical slots, use standard card orientation
                target_aspect = card_aspect  # 0.714
                
                # Calculate size to fit the vertical slot
                slot_aspect = slot_width_inches / slot_height_inches
                
                if target_aspect > slot_aspect:
                    # Image is wider than slot - fit to width
                    new_width = slot_width_inches
                    new_height = slot_width_inches / target_aspect
                else:
                    # Image is taller than slot - fit to height
                    new_height = slot_height_inches  
                    new_width = slot_height_inches * target_aspect
            
            # Convert to pixels (assuming 96 DPI)
            new_width_px = int(new_width * 96)
            new_height_px = int(new_height * 96)
            
            # Resize the image
            img_resized = img.resize((new_width_px, new_height_px), Image.Resampling.LANCZOS)
            
            # Save to bytes
            img_stream = io.BytesIO()
            img_resized.save(img_stream, format='PNG')
            img_stream.seek(0)
            
            return img_stream, new_width, new_height
            
    except Exception as e:
        print(f"    ❌ Error resizing {card_path}: {e}")
        return None, 0, 0

def place_cards_in_template_orientation_aware(template_file, card_images, output_file):
    """
    Place cards in the template respecting orientation requirements:
    - Slide 1: 2 vertical cards
    - Slide 2: 6 horizontal cards
    """
    print(f"🔧 Opening template: {template_file}")
    prs = Presentation(template_file)
    
    card_index = 0
    total_slots_found = 0
    
    # Process slides that already exist in template
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n📋 Processing Slide {slide_num} ({len(slide.shapes)} shapes)")
        
        # Find all card slots first
        card_slots = []
        for shape_num, shape in enumerate(slide.shapes):
            width_inches = shape.width / 914400
            height_inches = shape.height / 914400
            aspect_ratio = width_inches / height_inches if height_inches > 0 else 0
            
            # Detect Magic card slots by size and aspect ratio
            # Allow for both orientations in detection
            is_card_slot = False
            
            # Check for vertical card slot (2.5" × 3.5" approximate)
            if 0.65 < aspect_ratio < 0.8 and width_inches > 2.0 and height_inches > 3.0:
                is_card_slot = True
                
            # Check for horizontal card slot (3.5" × 2.5" approximate)  
            elif 1.25 < aspect_ratio < 1.55 and width_inches > 3.0 and height_inches > 2.0:
                is_card_slot = True
            
            if is_card_slot:
                orientation = analyze_card_slot_orientation(shape, slide_num)
                card_slots.append({
                    'shape': shape,
                    'shape_num': shape_num + 1,
                    'width_inches': width_inches,
                    'height_inches': height_inches,
                    'orientation': orientation
                })
        
        print(f"  🎴 Found {len(card_slots)} card slots on slide {slide_num}")
        
        # Replace each slot with a card image
        for slot_info in card_slots:
            if card_index >= len(card_images):
                print(f"  ⚠️  Ran out of cards at slot {total_slots_found + 1}")
                break
                
            card_path = card_images[card_index]
            card_name = os.path.basename(card_path)
            
            print(f"  🎴 Slot {total_slots_found + 1}: {card_name} ({slot_info['orientation']})")
            
            # Resize card for this specific slot with orientation
            img_stream, img_width, img_height = resize_card_for_slot(
                card_path, 
                slot_info['width_inches'], 
                slot_info['height_inches'],
                slot_info['orientation']
            )
            
            if img_stream:
                try:
                    # Remove the original shape
                    shape_element = slot_info['shape'].element
                    shape_element.getparent().remove(shape_element)
                    
                    # Add the card image in the same position
                    left = slot_info['shape'].left
                    top = slot_info['shape'].top
                    
                    # Add picture with proper size
                    new_shape = slide.shapes.add_picture(
                        img_stream, 
                        left, 
                        top, 
                        Inches(img_width),
                        Inches(img_height)
                    )
                    
                    print(f"    ✅ Placed {card_name} as {slot_info['orientation']} card")
                    
                except Exception as e:
                    print(f"    ❌ Error placing {card_name}: {e}")
            
            card_index += 1
            total_slots_found += 1
    
    # Add overflow cards in grid slides if needed
    remaining_cards = len(card_images) - card_index
    if remaining_cards > 0:
        print(f"\n📄 Creating grid slides for {remaining_cards} remaining cards...")
        
        cards_per_slide = 20
        slides_needed = (remaining_cards + cards_per_slide - 1) // cards_per_slide
        
        for grid_slide_num in range(slides_needed):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            
            start_idx = card_index
            end_idx = min(card_index + cards_per_slide, len(card_images))
            slide_cards = card_images[start_idx:end_idx]
            
            print(f"  📋 Grid slide {grid_slide_num + 1}: Cards {start_idx + 1}-{end_idx}")
            
            # Create 4x5 grid (standard orientation for overflow)
            cols, rows = 4, 5
            margin = 0.5
            card_width = (10 - 2 * margin - (cols - 1) * 0.2) / cols
            card_height = card_width / 0.714  # Standard Magic card aspect
            
            for i, card_path in enumerate(slide_cards):
                row = i // cols
                col = i % cols
                
                left = margin + col * (card_width + 0.2)
                top = margin + row * (card_height + 0.2)
                
                try:
                    with Image.open(card_path) as img:
                        # Resize for grid (vertical orientation)
                        width_px = int(card_width * 96)
                        height_px = int(card_height * 96)
                        img_resized = img.resize((width_px, height_px), Image.Resampling.LANCZOS)
                        
                        img_stream = io.BytesIO()
                        img_resized.save(img_stream, format='PNG')
                        img_stream.seek(0)
                        
                        slide.shapes.add_picture(
                            img_stream,
                            Inches(left),
                            Inches(top), 
                            Inches(card_width),
                            Inches(card_height)
                        )
                        
                except Exception as e:
                    print(f"    ❌ Error adding {os.path.basename(card_path)}: {e}")
            
            card_index += len(slide_cards)
    
    # Save the presentation
    print(f"\n💾 Saving to: {output_file}")
    prs.save(output_file)
    
    print(f"\n🎉 Template processing complete!")
    print(f"  📊 Total slides: {len(prs.slides)}")
    print(f"  🎴 Cards placed in template: {min(total_slots_found, len(card_images))}")
    print(f"  📄 Cards in grid slides: {max(0, len(card_images) - total_slots_found)}")
    print(f"  🎯 Total cards processed: {len(card_images)}")

def main():
    # File paths
    template_file = "magic_cards_template.pptx"
    output_file = "magic_cards_completed_ORIENTATION_FIXED.pptx"
    images_dir = "images"
    
    # Check if template exists
    if not os.path.exists(template_file):
        print(f"❌ Template file not found: {template_file}")
        return
    
    # Check if images directory exists
    if not os.path.exists(images_dir):
        print(f"❌ Images directory not found: {images_dir}")
        return
    
    # Get all card images (sorted by filename)
    card_images = sorted(glob.glob(os.path.join(images_dir, "*.jpg")))
    if not card_images:
        card_images = sorted(glob.glob(os.path.join(images_dir, "*.png")))
    
    if not card_images:
        print(f"❌ No card images found in {images_dir}")
        return
    
    print(f"🎴 Found {len(card_images)} card images")
    print(f"📁 Template: {template_file}")
    print(f"💾 Output: {output_file}")
    print(f"🎯 Processing orientation-aware placement...")
    
    # Process the template with orientation awareness
    place_cards_in_template_orientation_aware(template_file, card_images, output_file)

if __name__ == "__main__":
    main()