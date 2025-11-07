#!/usr/bin/env python3
import os
import json
import math
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image

def get_template_slot_positions(template_slide):
    """Extract card slot positions from template"""
    card_slots = []

    for shape in template_slide.shapes:
        if (shape.shape_type in [1, 13] and
            shape.width.inches > 1.0 and
            shape.height.inches > 1.0):

            card_slots.append({
                'left': shape.left.inches,
                'top': shape.top.inches,
                'width': shape.width.inches,
                'height': shape.height.inches
            })

    card_slots.sort(key=lambda x: (x['top'], x['left']))
    return card_slots

def place_card_in_slot(slide, card_path, slot_info):
    """Place card image in slot with aspect ratio preservation"""
    try:
        if not os.path.exists(card_path):
            print(f"    ⚠️  Image not found: {card_path}")
            return False

        with Image.open(card_path) as img:
            slot_aspect = slot_info['width'] / slot_info['height']
            img_aspect = img.width / img.height

            is_horizontal_slot = slot_aspect > 1.0

            if img_aspect > slot_aspect:
                new_width = slot_info['width']
                new_height = slot_info['width'] / img_aspect
            else:
                new_height = slot_info['height']
                new_width = slot_info['height'] * img_aspect

            left = slot_info['left'] + (slot_info['width'] - new_width) / 2
            top = slot_info['top'] + (slot_info['height'] - new_height) / 2

            target_width = int(new_width * 150)
            target_height = int(new_height * 150)

            resized_img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)

            if is_horizontal_slot and img_aspect < 1.0:
                resized_img = resized_img.rotate(90, expand=True)

            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as temp_file:
                temp_path = temp_file.name
                resized_img.save(temp_path, "JPEG", quality=90)

        slide.shapes.add_picture(
            temp_path,
            Inches(left),
            Inches(top),
            Inches(new_width),
            Inches(new_height)
        )

        os.remove(temp_path)
        return True

    except Exception as e:
        print(f"    ❌ Failed to place card: {e}")
        return False

if __name__ == "__main__":
    template_file = "decks.pptx"
    output_file = "krenko_mob_boss.pptx"
    manifest_path = ".claude/state/krenko_manifest.json"

    print(f"🎴 Generating Krenko Mob Boss Presentation")
    print("=" * 70)

    # Load manifest
    with open(manifest_path, 'r') as f:
        manifest = json.load(f)

    # Flatten cards for DFC support
    card_images = []
    for card in manifest['cards']:
        if card['status'] != 'success':
            continue

        if 'paths' in card and card.get('faces', 0) > 1:
            for i, face_path in enumerate(card['paths']):
                card_images.append({
                    'name': f"{card['name']} (Face {i+1})",
                    'path': face_path
                })
        else:
            card_images.append({
                'name': card['name'],
                'path': card['path']
            })

    total_cards = len(card_images)
    print(f"🎴 Found {total_cards} card images")

    # Load template
    template_prs = Presentation(template_file)
    template_slide = template_prs.slides[0]
    slot_positions = get_template_slot_positions(template_slide)

    print(f"📐 Found {len(slot_positions)} card slots in template")

    cards_per_slide = len(slot_positions)
    slides_needed = math.ceil(total_cards / cards_per_slide)

    print(f"📄 Creating {slides_needed} slides\n")

    # Create presentation from template
    prs = Presentation(template_file)

    # Remove template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    blank_layout = prs.slide_layouts[6]

    # Generate slides
    for slide_num in range(slides_needed):
        slide = prs.slides.add_slide(blank_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        start_idx = slide_num * cards_per_slide
        end_idx = min(start_idx + cards_per_slide, total_cards)

        print(f"📄 Slide {slide_num + 1}: Cards {start_idx + 1}-{end_idx}")

        for i, card_idx in enumerate(range(start_idx, end_idx)):
            if i >= len(slot_positions):
                break

            card_info = card_images[card_idx]
            place_card_in_slot(slide, card_info['path'], slot_positions[i])

    prs.save(output_file)

    # Verify dimensions
    template = Presentation(template_file)
    output = Presentation(output_file)

    print(f"\n💾 Saved: {output_file}")
    print(f"📊 Stats: {total_cards} cards across {slides_needed} slides")
    print(f"📐 Dimensions: {output.slide_width.inches:.1f}\" x {output.slide_height.inches:.1f}\"")
    print(f"✅ Printer ready!")
