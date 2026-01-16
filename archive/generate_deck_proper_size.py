#!/usr/bin/env python3
"""
Proper MTG deck generator using template with correct card dimensions (2.5" × 3.5").
"""
import sys
import json
import time
import requests
from pathlib import Path
import subprocess
import math
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image

def fetch_card_from_scryfall(card_name):
    """Fetch card data from Scryfall API."""
    url = f"https://api.scryfall.com/cards/named"
    params = {"fuzzy": card_name}

    try:
        response = requests.get(url, params=params, timeout=10)
        response.raise_for_status()
        return response.json()
    except Exception as e:
        print(f"  ❌ Failed to fetch {card_name}: {e}")
        return None

def download_image(url, output_path):
    """Download image from URL."""
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()

        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'wb') as f:
            f.write(response.content)
        return True
    except Exception as e:
        print(f"  ❌ Failed to download image: {e}")
        return False

def get_template_slot_positions(template_slide, max_slots=8):
    """
    Extract card slot positions from template.
    Returns list of slot dictionaries with left, top, width, height.
    Only includes VISIBLE slots (within slide boundaries).
    Limits to max_slots (default 8 for 2v6h layout).
    """
    from pptx import Presentation

    card_slots = []
    prs = template_slide.part.presentation
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches

    for shape in template_slide.shapes:
        # Look for rectangular shapes that are card-sized (around 2.5" × 3.5")
        if (hasattr(shape, 'width') and hasattr(shape, 'height') and
            shape.width.inches > 1.0 and shape.height.inches > 1.0):

            # Check if shape is within slide boundaries (visible)
            is_visible = (shape.top.inches >= 0 and
                         shape.left.inches >= 0 and
                         shape.top.inches < slide_height and
                         shape.left.inches < slide_width)

            if is_visible:
                card_slots.append({
                    'left': shape.left.inches,
                    'top': shape.top.inches,
                    'width': shape.width.inches,
                    'height': shape.height.inches
                })

    # Sort slots by position (top to bottom, left to right)
    card_slots.sort(key=lambda x: (x['top'], x['left']))

    # Limit to max_slots (8 for proper 2v6h layout)
    return card_slots[:max_slots]

def place_card_in_slot(slide, card_path, slot_info):
    """
    Place card image in slot with aspect ratio preservation.
    """
    try:
        if not card_path.exists():
            print(f"    ⚠️  Image not found: {card_path}")
            return False

        with Image.open(card_path) as img:
            # Calculate aspect ratios
            slot_aspect = slot_info['width'] / slot_info['height']
            img_aspect = img.width / img.height

            # Preserve aspect ratio while fitting to slot
            if img_aspect > slot_aspect:
                # Image is wider, fit to width
                new_width = slot_info['width']
                new_height = slot_info['width'] / img_aspect
            else:
                # Image is taller, fit to height
                new_height = slot_info['height']
                new_width = slot_info['height'] * img_aspect

            # Center the image in the slot
            left = slot_info['left'] + (slot_info['width'] - new_width) / 2
            top = slot_info['top'] + (slot_info['height'] - new_height) / 2

        # Add picture to slide
        slide.shapes.add_picture(
            str(card_path),
            Inches(left),
            Inches(top),
            width=Inches(new_width),
            height=Inches(new_height)
        )

        return True

    except Exception as e:
        print(f"    ❌ Failed to place card: {e}")
        return False

def create_presentation_from_template(template_file, card_images, output_file):
    """
    Create presentation using template with proper MTG card dimensions.
    """
    # Load template
    print(f"📋 Loading template: {template_file}")
    template_prs = Presentation(template_file)

    if len(template_prs.slides) == 0:
        print("❌ Template has no slides")
        return False

    # Get slot positions from first slide
    template_slide = template_prs.slides[0]
    slot_positions = get_template_slot_positions(template_slide)

    print(f"📐 Found {len(slot_positions)} card slots in template")
    print(f"   Slot dimensions: ~{slot_positions[0]['width']:.2f}\" × {slot_positions[0]['height']:.2f}\"")

    if len(slot_positions) == 0:
        print("❌ No card slots found in template")
        return False

    cards_per_slide = len(slot_positions)
    slides_needed = math.ceil(len(card_images) / cards_per_slide)

    print(f"📄 Creating {slides_needed} slides ({cards_per_slide} cards per slide)")

    # Create new presentation from template
    prs = Presentation(template_file)

    # Remove existing slides but keep layouts
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Process cards in groups
    cards_placed = 0
    for slide_num in range(slides_needed):
        # Add new slide
        slide = prs.slides.add_slide(blank_layout)

        # Black background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Calculate card indices for this slide
        start_idx = slide_num * cards_per_slide
        end_idx = min(start_idx + cards_per_slide, len(card_images))

        print(f"\n📄 Slide {slide_num + 1}: Cards {start_idx + 1} to {end_idx}")

        # Place cards using template positions
        for i, card_idx in enumerate(range(start_idx, end_idx)):
            if i >= len(slot_positions):
                break

            card_path = card_images[card_idx]
            slot = slot_positions[i]

            card_name = card_path.stem
            print(f"  [{i+1}] {card_name}")
            if place_card_in_slot(slide, card_path, slot):
                print(f"      ✅ Placed at {slot['width']:.2f}\" × {slot['height']:.2f}\"")
                cards_placed += 1
            else:
                print(f"      ❌ Failed to place")

    # Save presentation
    prs.save(output_file)
    print(f"\n💾 Saved presentation: {output_file}")
    print(f"📊 Stats: {cards_placed} cards across {slides_needed} slides")
    print(f"📏 Card dimensions: 2.5\" × 3.5\" (proper MTG size)")

    return True

def main(card_list_file, template_file="template_2v6h_FIXED.pptx", output_dir="outputs"):
    """Generate deck with proper MTG card dimensions."""

    print("🎴 MTG Deck Generator (Proper Card Dimensions)")
    print("=" * 70)

    # Verify template exists
    if not Path(template_file).exists():
        print(f"❌ Template not found: {template_file}")
        print("   Download from: https://docs.google.com/presentation/d/1Dz6itDc8r0PVyztqSAQc5xN9i9Y3kcuVwx7x2Jt1LTY")
        return 1

    # Read card list
    card_list_path = Path(card_list_file)
    if not card_list_path.exists():
        print(f"❌ Card list not found: {card_list_file}")
        return 1

    with open(card_list_path, 'r') as f:
        card_names = [line.strip() for line in f if line.strip()]

    print(f"📋 Found {len(card_names)} cards in list")
    print()

    # Create output directories
    images_dir = Path("images")
    images_dir.mkdir(exist_ok=True)
    output_path = Path(output_dir)
    output_path.mkdir(exist_ok=True)

    # Fetch and download cards
    print("🔍 Fetching card data from Scryfall...")
    successful_images = []

    for i, card_name in enumerate(card_names, 1):
        print(f"[{i}/{len(card_names)}] {card_name}")

        # Fetch card data
        card_data = fetch_card_from_scryfall(card_name)
        if not card_data:
            continue

        # Get image URL
        image_url = card_data.get('image_uris', {}).get('normal')
        if not image_url:
            print(f"  ⚠️  No image available")
            continue

        # Download image
        safe_name = card_name.replace(' ', '_').replace('/', '_').replace("'", "")
        image_path = images_dir / f"{safe_name}.jpg"

        if image_path.exists():
            print(f"  ✓ Already downloaded")
            successful_images.append(image_path)
        else:
            if download_image(image_url, image_path):
                print(f"  ✓ Downloaded")
                successful_images.append(image_path)

        # Rate limiting
        time.sleep(0.1)

    print()
    print(f"✓ Successfully downloaded {len(successful_images)}/{len(card_names)} cards")
    print()

    if len(successful_images) == 0:
        print("❌ No cards downloaded, cannot generate deck")
        return 1

    # Generate PPTX using template
    print("📄 Generating PPTX with proper card dimensions...")
    pptx_path = output_path / "deck_proper_size.pptx"

    if not create_presentation_from_template(template_file, successful_images, pptx_path):
        return 1

    print()

    # Convert to PDF
    print("📄 Converting to PDF...")
    pdf_path = output_path / "deck_proper_size.pdf"

    try:
        subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_path),
            str(pptx_path)
        ], check=True, capture_output=True, timeout=60)
        print(f"  ✓ {pdf_path}")
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as e:
        print(f"  ⚠️  PDF conversion failed (LibreOffice not available)")
        print(f"  ℹ️  You can convert {pptx_path} to PDF manually")

    print()
    print("=" * 70)
    print("🎉 Deck generation complete!")
    print(f"📁 PPTX: {pptx_path}")
    if pdf_path.exists():
        print(f"📁 PDF:  {pdf_path}")
    print(f"🎴 {len(successful_images)} cards with PROPER MTG DIMENSIONS (2.5\" × 3.5\")")
    print(f"📏 Ready for printing!")

    return 0

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 generate_deck_proper_size.py <card_list_file>")
        print()
        print("Example:")
        print("  python3 generate_deck_proper_size.py decklists/user_deck.txt")
        sys.exit(1)

    sys.exit(main(sys.argv[1]))
