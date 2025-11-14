#!/usr/bin/env python3
import os
import json
import math
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image
from pathlib import Path

def get_template_slot_positions(template_slide):
    """Extract card slot positions from template"""
    card_slots = []

    for shape in template_slide.shapes:
        if (shape.shape_type in [1, 13] and
            shape.width.inches > 1.0 and
            shape.height.inches > 1.0):

            card_slots.append({
                'left': shape.left.inches,
                'top': shape.top.inches,
                'width': shape.width.inches,
                'height': shape.height.inches
            })

    card_slots.sort(key=lambda x: (x['top'], x['left']))
    return card_slots

def place_card_in_slot(slide, card_path, slot_info):
    """Place card image in slot with aspect ratio preservation"""
    try:
        if not os.path.exists(card_path):
            print(f"    ⚠️  Image not found: {card_path}")
            return False

        with Image.open(card_path) as img:
            slot_aspect = slot_info['width'] / slot_info['height']
            img_aspect = img.width / img.height

            is_horizontal_slot = slot_aspect > 1.0

            if img_aspect > slot_aspect:
                new_width = slot_info['width']
                new_height = slot_info['width'] / img_aspect
            else:
                new_height = slot_info['height']
                new_width = slot_info['height'] * img_aspect

            left = slot_info['left'] + (slot_info['width'] - new_width) / 2
            top = slot_info['top'] + (slot_info['height'] - new_height) / 2

            target_width = int(new_width * 150)
            target_height = int(new_height * 150)

            resized_img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)

            if is_horizontal_slot and img_aspect < 1.0:
                resized_img = resized_img.rotate(90, expand=True)

            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as temp_file:
                temp_path = temp_file.name
                resized_img.save(temp_path, "JPEG", quality=90)

        slide.shapes.add_picture(
            temp_path,
            Inches(left),
            Inches(top),
            Inches(new_width),
            Inches(new_height)
        )

        os.remove(temp_path)
        return True

    except Exception as e:
        print(f"    ❌ Failed to place card: {e}")
        return False

def create_presentation_from_template(template_file, images_dir, output_file, manifest_path):
    """Create presentation using template pattern"""

    if not os.path.exists(template_file):
        print(f"❌ Error: Template file not found: {template_file}")
        return False

    if not os.path.exists(manifest_path):
        print(f"❌ Error: Manifest not found: {manifest_path}")
        return False

    with open(manifest_path, 'r') as f:
        manifest = json.load(f)

    successful_cards = [
        card for card in manifest['cards']
        if card['status'] == 'success'
    ]

    total_cards = len(successful_cards)
    print(f"🎴 Found {total_cards} cards to place from manifest")

    if total_cards == 0:
        print("❌ No successful cards in manifest")
        return False

    print(f"📋 Loading template: {template_file}")
    template_prs = Presentation(template_file)

    if len(template_prs.slides) == 0:
        print("❌ Template has no slides")
        return False

    template_slide = template_prs.slides[0]
    slot_positions = get_template_slot_positions(template_slide)

    print(f"📐 Found {len(slot_positions)} card slots in template")

    if len(slot_positions) == 0:
        print("❌ No card slots found in template")
        return False

    cards_per_slide = len(slot_positions)
    slides_needed = math.ceil(total_cards / cards_per_slide)

    print(f"📄 Creating {slides_needed} slides ({cards_per_slide} cards per slide)")

    # Create new presentation from template to preserve dimensions
    prs = Presentation(template_file)

    # Remove template slides, keep only the slide master/layouts
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    blank_layout = prs.slide_layouts[6]  # Blank layout from template

    for slide_num in range(slides_needed):
        slide = prs.slides.add_slide(blank_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        start_idx = slide_num * cards_per_slide
        end_idx = min(start_idx + cards_per_slide, total_cards)

        print(f"\n📄 Slide {slide_num + 1}: Cards {start_idx + 1} to {end_idx}")

        for i, card_idx in enumerate(range(start_idx, end_idx)):
            if i >= len(slot_positions):
                break

            card_info = successful_cards[card_idx]
            card_path = card_info['path']
            card_name = card_info['name']
            slot = slot_positions[i]

            print(f"  [{i+1}] {card_name}")
            if place_card_in_slot(slide, card_path, slot):
                print(f"      ✅ Placed in slot")
            else:
                print(f"      ❌ Failed to place")

        remaining_slots = len(slot_positions) - (end_idx - start_idx)
        if remaining_slots > 0:
            print(f"  ℹ️  {remaining_slots} empty slots left blank")

    prs.save(output_file)
    print(f"\n💾 Saved presentation: {output_file}")
    print(f"📊 Stats: {total_cards} cards across {slides_needed} slides")

    return True

if __name__ == "__main__":
    print("🎴 Regenerating Phase 1 Output with Correct Dimensions")
    print("=" * 70)

    success = create_presentation_from_template(
        'decks.pptx',
        'images',
        'tests/phase1_output.pptx',
        '.claude/state/fetch_manifest.json'
    )

    if success:
        print("\n" + "=" * 70)
        print("🎉 Phase 1 output regenerated successfully!")

        # Verify dimensions
        from pptx import Presentation
        template = Presentation('decks.pptx')
        output = Presentation('tests/phase1_output.pptx')

        print(f"\nTemplate: {template.slide_width.inches:.1f}\" x {template.slide_height.inches:.1f}\"")
        print(f"Output:   {output.slide_width.inches:.1f}\" x {output.slide_height.inches:.1f}\"")

        if template.slide_width == output.slide_width and template.slide_height == output.slide_height:
            print("✅ Dimensions match - printer ready!")
