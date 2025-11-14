#!/usr/bin/env python3
"""
Detailed Template Analysis Script
Analyzes the actual PowerPoint template structure to understand layout
"""

from pptx import Presentation
import os

def analyze_template_structure(template_file):
    """Analyze the template structure in detail"""
    print(f"🔍 Analyzing template: {template_file}")
    
    if not os.path.exists(template_file):
        print(f"❌ Template not found: {template_file}")
        return
    
    prs = Presentation(template_file)
    
    print(f"📊 Template has {len(prs.slides)} slides")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n🔸 Slide {slide_num}:")
        print(f"  📋 Total shapes: {len(slide.shapes)}")
        
        card_slots = []
        
        for shape_num, shape in enumerate(slide.shapes):
            width_inches = shape.width / 914400
            height_inches = shape.height / 914400
            left_inches = shape.left / 914400
            top_inches = shape.top / 914400
            aspect_ratio = width_inches / height_inches if height_inches > 0 else 0
            
            print(f"    Shape {shape_num + 1}: {width_inches:.2f}×{height_inches:.2f} at ({left_inches:.2f},{top_inches:.2f}) AR:{aspect_ratio:.3f}")
            
            # Check if this looks like a Magic card slot
            # Magic cards are 2.5" × 3.5" (AR ≈ 0.714) or rotated 3.5" × 2.5" (AR ≈ 1.4)
            is_vertical_card = 0.65 < aspect_ratio < 0.8 and width_inches > 2.0 and height_inches > 3.0
            is_horizontal_card = 1.25 < aspect_ratio < 1.55 and width_inches > 3.0 and height_inches > 2.0
            
            if is_vertical_card:
                card_slots.append({
                    'shape_num': shape_num + 1,
                    'orientation': 'vertical',
                    'size': f"{width_inches:.2f}×{height_inches:.2f}",
                    'position': f"({left_inches:.2f},{top_inches:.2f})",
                    'aspect_ratio': aspect_ratio
                })
                print(f"      🎴 VERTICAL card slot detected")
                
            elif is_horizontal_card:
                card_slots.append({
                    'shape_num': shape_num + 1,
                    'orientation': 'horizontal', 
                    'size': f"{width_inches:.2f}×{height_inches:.2f}",
                    'position': f"({left_inches:.2f},{top_inches:.2f})",
                    'aspect_ratio': aspect_ratio
                })
                print(f"      🎴 HORIZONTAL card slot detected")
        
        print(f"\n  📊 Card slots found on slide {slide_num}: {len(card_slots)}")
        for slot in card_slots:
            print(f"    🎴 Slot {slot['shape_num']}: {slot['orientation']} {slot['size']} at {slot['position']}")
        
        # Analyze layout patterns
        if len(card_slots) > 0:
            vertical_count = sum(1 for slot in card_slots if slot['orientation'] == 'vertical')
            horizontal_count = sum(1 for slot in card_slots if slot['orientation'] == 'horizontal')
            
            print(f"  📈 Layout summary: {vertical_count} vertical, {horizontal_count} horizontal")
            
            # Group by position to understand layout
            positions = [(float(slot['position'].split(',')[0][1:]), float(slot['position'].split(',')[1][:-1])) for slot in card_slots]
            
            # Sort by left position to see columns
            positions_sorted = sorted(positions)
            print(f"  📍 Positions (left to right): {positions_sorted}")

def main():
    template_file = "magic_cards_template.pptx"
    analyze_template_structure(template_file)

if __name__ == "__main__":
    main()