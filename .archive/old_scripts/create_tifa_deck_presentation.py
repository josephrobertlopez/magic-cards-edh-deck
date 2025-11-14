#!/usr/bin/env python3
"""
Create Tifa Lockhart EDH Deck Presentation - 75 Cards
Using PNG template layout: 2 vertical slots + 6 horizontal slots = 8 cards per page
"""

import os
import glob
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import math

def resize_card_for_slot(card_path, slot_width_inches, slot_height_inches, orientation='vertical'):
    """
    Resize a Magic card image to fit a slot with specific orientation
    """
    try:
        with Image.open(card_path) as img:
            # Standard Magic card aspect ratio is 2.5:3.5 (width:height) = 0.714
            card_aspect = 2.5 / 3.5
            
            if orientation == 'horizontal':
                # For horizontal slots, rotate the card aspect ratio
                card_aspect = 3.5 / 2.5  # 1.4
                
                # Rotate the image 90 degrees for horizontal placement
                img = img.rotate(90, expand=True)
            
            # Calculate size to fit the slot
            slot_aspect = slot_width_inches / slot_height_inches
            
            if card_aspect > slot_aspect:
                # Image is wider than slot - fit to width
                new_width = slot_width_inches
                new_height = slot_width_inches / card_aspect
            else:
                # Image is taller than slot - fit to height
                new_height = slot_height_inches
                new_width = slot_height_inches * card_aspect
            
            # Convert to pixels (assuming 96 DPI)
            new_width_px = int(new_width * 96)
            new_height_px = int(new_height * 96)
            
            # Resize the image
            img_resized = img.resize((new_width_px, new_height_px), Image.Resampling.LANCZOS)
            
            # Save to bytes
            img_stream = io.BytesIO()
            img_resized.save(img_stream, format='PNG')
            img_stream.seek(0)
            
            return img_stream, new_width, new_height
            
    except Exception as e:
        print(f"    ❌ Error resizing {card_path}: {e}")
        return None, 0, 0

def create_tifa_deck_presentation(template_file, card_images, output_file):
    """
    Create Tifa Lockhart EDH deck presentation with 75 cards
    Layout: 2 vertical slots (left) + 6 horizontal slots (right, 2x3 grid)
    """
    print(f"🔧 Opening template: {template_file}")
    prs = Presentation(template_file)
    
    if len(prs.slides) < 1:
        print("❌ Template has no slides!")
        return
    
    total_cards = len(card_images)
    cards_per_slide = 8
    slides_needed = math.ceil(total_cards / cards_per_slide)
    
    print(f"🎴 Total cards in Tifa's deck: {total_cards}")
    print(f"📄 Slides needed: {slides_needed} (8 cards per slide)")
    
    # Define slot positions based on PNG template layout
    slot_positions = [
        # 2 vertical slots (left column)
        {'x': 0.5, 'y': 1.0, 'width': 2.5, 'height': 3.5, 'orientation': 'vertical'},   # Top left
        {'x': 0.5, 'y': 5.0, 'width': 2.5, 'height': 3.5, 'orientation': 'vertical'},   # Bottom left
        
        # 6 horizontal slots (right side, 2x3 grid)
        {'x': 3.5, 'y': 0.5, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'}, # Top row, left
        {'x': 7.5, 'y': 0.5, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'}, # Top row, right
        {'x': 3.5, 'y': 3.5, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'}, # Middle row, left
        {'x': 7.5, 'y': 3.5, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'}, # Middle row, right
        {'x': 3.5, 'y': 6.5, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'}, # Bottom row, left
        {'x': 7.5, 'y': 6.5, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'}, # Bottom row, right
    ]
    
    # Process cards slide by slide
    for slide_num in range(slides_needed):
        # Create new slide if needed (keep slide 1, create additional)
        if slide_num == 0:
            # Use existing first slide but clear it
            slide = prs.slides[0]
            # Remove all existing shapes
            for shape in list(slide.shapes):
                if hasattr(shape.element, 'getparent'):
                    shape.element.getparent().remove(shape.element)
        else:
            # Add new slide using same layout
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
        
        print(f"\n📄 Processing Slide {slide_num + 1}")
        
        # Place cards on this slide
        start_card = slide_num * cards_per_slide
        end_card = min(start_card + cards_per_slide, total_cards)
        slide_cards = card_images[start_card:end_card]
        
        print(f"  🎴 Cards {start_card + 1}-{end_card} ({len(slide_cards)} cards)")
        
        # Place each card in its designated slot
        for i, card_path in enumerate(slide_cards):
            if i >= len(slot_positions):
                break
                
            slot = slot_positions[i]
            card_name = os.path.basename(card_path).replace('.jpg', '')
            
            print(f"    🎴 Slot {i + 1}: {card_name} ({slot['orientation']})")
            
            # Resize card for this specific slot
            img_stream, img_width, img_height = resize_card_for_slot(
                card_path, 
                slot['width'], 
                slot['height'],
                slot['orientation']
            )
            
            if img_stream:
                try:
                    # Add the card image
                    left = Inches(slot['x'])
                    top = Inches(slot['y'])
                    
                    new_shape = slide.shapes.add_picture(
                        img_stream, 
                        left, 
                        top, 
                        Inches(img_width),
                        Inches(img_height)
                    )
                    
                    print(f"      ✅ Placed {card_name}")
                    
                except Exception as e:
                    print(f"      ❌ Error placing {card_name}: {e}")
    
    # Save the presentation
    print(f"\n💾 Saving to: {output_file}")
    prs.save(output_file)
    
    print(f"\n🎉 Tifa Lockhart EDH deck presentation complete!")
    print(f"  📊 Total slides: {len(prs.slides)}")
    print(f"  🎴 Total cards placed: {total_cards}")
    print(f"  📋 Cards per slide: {cards_per_slide} (2 vertical + 6 horizontal)")
    print(f"  🎮 Commander: Tifa Lockhart")

def main():
    # File paths
    template_file = "magic_cards_template.pptx"
    output_file = "Tifa_Lockhart_EDH_Deck.pptx"
    images_dir = "images"
    
    # Check if template exists
    if not os.path.exists(template_file):
        print(f"❌ Template file not found: {template_file}")
        return
    
    # Check if images directory exists
    if not os.path.exists(images_dir):
        print(f"❌ Images directory not found: {images_dir}")
        return
    
    # Get all card images (sorted by filename for consistent ordering)
    card_images = sorted(glob.glob(os.path.join(images_dir, "*.jpg")))
    
    if not card_images:
        print(f"❌ No card images found in {images_dir}")
        return
    
    print(f"🎴 Found {len(card_images)} cards for Tifa Lockhart's EDH deck")
    print(f"📁 Template: {template_file}")
    print(f"💾 Output: {output_file}")
    print(f"🎯 Creating Tifa Lockhart EDH deck presentation...")
    
    # Create the presentation
    create_tifa_deck_presentation(template_file, card_images, output_file)

if __name__ == "__main__":
    main()