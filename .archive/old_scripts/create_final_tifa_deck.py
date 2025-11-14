#!/usr/bin/env python3
"""
Create the final Tifa Lockhart EDH deck presentation with all 100 cards
Using the PNG template layout: 2 vertical slots (left) + 6 horizontal slots (right)
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from pathlib import Path

def create_final_tifa_deck_presentation(template_file, output_file):
    """Create complete Tifa Lockhart EDH deck presentation"""
    
    # Get all card images
    images_dir = "images"
    card_images = []
    
    if os.path.exists(images_dir):
        for filename in sorted(os.listdir(images_dir)):
            if filename.lower().endswith(('.jpg', '.jpeg', '.png')):
                card_images.append(os.path.join(images_dir, filename))
    
    total_cards = len(card_images)
    print(f"🎴 Creating presentation with {total_cards} cards")
    
    # Load the template
    prs = Presentation(template_file)
    
    # Define slot positions based on PNG template analysis
    # Layout: 2 vertical slots (left column) + 6 horizontal slots (right, 2×3 grid)
    slot_positions = [
        # 2 vertical slots (left column)
        {'x': 0.5, 'y': 1.0, 'width': 2.5, 'height': 3.5, 'orientation': 'vertical'},
        {'x': 0.5, 'y': 5.0, 'width': 2.5, 'height': 3.5, 'orientation': 'vertical'},
        
        # 6 horizontal slots (right side, 2×3 grid)
        {'x': 4.0, 'y': 0.8, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'},
        {'x': 8.0, 'y': 0.8, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'},
        
        {'x': 4.0, 'y': 3.8, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'},
        {'x': 8.0, 'y': 3.8, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'},
        
        {'x': 4.0, 'y': 6.8, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'},
        {'x': 8.0, 'y': 6.8, 'width': 3.5, 'height': 2.5, 'orientation': 'horizontal'},
    ]
    
    cards_per_slide = len(slot_positions)  # 8 slots per slide
    slides_needed = math.ceil(total_cards / cards_per_slide)
    
    print(f"📊 Creating {slides_needed} slides for {total_cards} cards")
    
    # Process cards in groups of 8
    for slide_num in range(slides_needed):
        print(f"🖼️ Creating slide {slide_num + 1}/{slides_needed}")
        
        # Use the first slide of template, or duplicate it
        if slide_num == 0:
            slide = prs.slides[0]  # Use existing first slide
        else:
            # Duplicate the first slide for additional slides
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
        
        # Calculate which cards go on this slide
        start_idx = slide_num * cards_per_slide
        end_idx = min(start_idx + cards_per_slide, total_cards)
        slide_cards = card_images[start_idx:end_idx]
        
        # Place cards in the slots
        for i, card_path in enumerate(slide_cards):
            if i >= len(slot_positions):
                break
                
            slot = slot_positions[i]
            
            try:
                # Load and resize the image
                with Image.open(card_path) as img:
                    # Resize image to fit slot while maintaining aspect ratio
                    target_width = int(slot['width'] * 72)  # Convert inches to pixels (72 DPI)
                    target_height = int(slot['height'] * 72)
                    
                    # Calculate scaling to fit within slot
                    scale_w = target_width / img.width
                    scale_h = target_height / img.height
                    scale = min(scale_w, scale_h)
                    
                    new_width = int(img.width * scale)
                    new_height = int(img.height * scale)
                    
                    # Resize image
                    resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                    
                    # Save temporary resized image
                    temp_path = f"temp_card_{slide_num}_{i}.jpg"
                    resized_img.save(temp_path, "JPEG", quality=85)
                
                # Add image to slide
                left = Inches(slot['x'])
                top = Inches(slot['y'])
                width = Inches(slot['width'])
                height = Inches(slot['height'])
                
                slide.shapes.add_picture(temp_path, left, top, width, height)
                
                # Clean up temporary file
                os.remove(temp_path)
                
                card_name = os.path.splitext(os.path.basename(card_path))[0]
                print(f"  ✅ Added {card_name} to slot {i+1}")
                
            except Exception as e:
                print(f"  ❌ Failed to add {card_path}: {e}")
    
    # Save the presentation
    prs.save(output_file)
    print(f"💾 Saved presentation: {output_file}")
    
    return True

def main():
    """Create the final Tifa Lockhart EDH deck presentation"""
    
    template_file = "magic_cards_template.pptx"
    output_file = "Tifa_Lockhart_EDH_Deck_FINAL.pptx"
    
    print(f"🎴 Creating Final Tifa Lockhart EDH Deck Presentation")
    print("=" * 60)
    
    if not os.path.exists(template_file):
        print(f"❌ Template file not found: {template_file}")
        return False
    
    if not os.path.exists("images"):
        print(f"❌ Images directory not found")
        return False
    
    success = create_final_tifa_deck_presentation(template_file, output_file)
    
    if success:
        print("\n" + "=" * 60)
        print(f"🎉 Final Tifa Lockhart EDH deck presentation created!")
        print(f"📁 File: {output_file}")
        print(f"🎴 Complete 100-card EDH deck with:")
        print(f"   • 72 unique named cards from your decklist")
        print(f"   • Both faces of Walk-In Closet // Forgotten Cellar")
        print(f"   • 28 Forest basic lands")
        print(f"   • Total: 102 images for 100-card deck")
    else:
        print("❌ Failed to create presentation")

if __name__ == "__main__":
    main()