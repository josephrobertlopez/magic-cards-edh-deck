#!/usr/bin/env python3
"""
Magic Cards PowerPoint Automation
Automatically places Magic card images into PowerPoint template slots
"""

import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import math

def analyze_presentation_template(pptx_file):
    """Analyze the PowerPoint template to find card slots"""
    print("🔍 Analyzing presentation template...")
    
    prs = Presentation(pptx_file)
    print(f"✅ Found {len(prs.slides)} slides in template")
    
    slide_info = []
    total_placeholders = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            'slide_number': slide_idx + 1,
            'shapes': [],
            'placeholders': []
        }
        
        print(f"\n📄 Slide {slide_idx + 1}:")
        
        # Analyze all shapes on the slide
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = {
                'index': shape_idx,
                'name': getattr(shape, 'name', 'Unknown'),
                'type': shape.shape_type,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }
            
            # Check if shape could be a card placeholder
            is_placeholder = False
            
            # Look for rectangular shapes that could hold cards
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or 
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE or
                shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER):
                
                # Check if it's roughly card-sized (Magic cards are ~2.5" x 3.5")
                width_inches = shape.width / 914400  # Convert EMU to inches
                height_inches = shape.height / 914400
                
                # Accept shapes between 1" and 5" in both dimensions
                if (1.0 <= width_inches <= 5.0 and 1.0 <= height_inches <= 5.0):
                    is_placeholder = True
                    shape_info['is_card_slot'] = True
                    shape_info['width_inches'] = width_inches
                    shape_info['height_inches'] = height_inches
                    slide_data['placeholders'].append(shape_info)
                    total_placeholders += 1
                    
                    print(f"   🎴 Card slot {total_placeholders}: {width_inches:.1f}\" x {height_inches:.1f}\"")
            
            if not is_placeholder:
                shape_info['is_card_slot'] = False
            
            slide_data['shapes'].append(shape_info)
        
        print(f"   📊 Found {len(slide_data['placeholders'])} card slots on this slide")
        slide_info.append(slide_data)
    
    print(f"\n🎯 Total card slots found: {total_placeholders}")
    return prs, slide_info, total_placeholders

def get_card_images(images_dir="images"):
    """Get all Magic card images"""
    print(f"\n🃏 Scanning {images_dir} for card images...")
    
    image_files = []
    for ext in ['*.jpg', '*.jpeg', '*.png']:
        image_files.extend(glob.glob(os.path.join(images_dir, ext)))
    
    image_files.sort()  # Sort for consistent ordering
    
    print(f"✅ Found {len(image_files)} card images")
    return image_files

def resize_image_for_slot(image_path, target_width_inches, target_height_inches, temp_dir="temp_resized"):
    """Resize image to fit card slot while maintaining aspect ratio"""
    
    # Create temp directory if it doesn't exist
    os.makedirs(temp_dir, exist_ok=True)
    
    # Open and resize image
    with Image.open(image_path) as img:
        # Calculate target size in pixels (assuming 96 DPI)
        target_width_px = int(target_width_inches * 96)
        target_height_px = int(target_height_inches * 96)
        
        # Resize while maintaining aspect ratio
        img.thumbnail((target_width_px, target_height_px), Image.Resampling.LANCZOS)
        
        # Save resized image
        filename = os.path.basename(image_path)
        temp_path = os.path.join(temp_dir, f"resized_{filename}")
        img.save(temp_path, "JPEG", quality=90)
        
        return temp_path

def place_cards_in_template(prs, slide_info, card_images):
    """Place Magic card images into the template slots"""
    print(f"\n🎨 Placing {len(card_images)} cards into template...")
    
    card_index = 0
    placed_cards = 0
    
    for slide_data in slide_info:
        slide_num = slide_data['slide_number']
        slide = prs.slides[slide_num - 1]
        placeholders = slide_data['placeholders']
        
        if not placeholders:
            continue
            
        print(f"\n📄 Processing slide {slide_num} ({len(placeholders)} slots)...")
        
        for placeholder in placeholders:
            if card_index >= len(card_images):
                print("⚠️ Ran out of card images!")
                break
            
            card_image = card_images[card_index]
            card_name = os.path.basename(card_image)
            
            try:
                # Resize image to fit slot
                resized_image = resize_image_for_slot(
                    card_image, 
                    placeholder['width_inches'], 
                    placeholder['height_inches']
                )
                
                # Add image to slide at placeholder position
                left = placeholder['left']
                top = placeholder['top']
                width = placeholder['width']
                height = placeholder['height']
                
                picture = slide.shapes.add_picture(resized_image, left, top, width, height)
                
                print(f"   ✅ Placed {card_name} in slot {placed_cards + 1}")
                placed_cards += 1
                card_index += 1
                
            except Exception as e:
                print(f"   ❌ Failed to place {card_name}: {e}")
                card_index += 1
                continue
    
    print(f"\n🎉 Successfully placed {placed_cards} cards!")
    return placed_cards

def create_card_grid_slides(prs, card_images, cards_per_slide=20):
    """Create additional slides with card grids if needed"""
    print(f"\n📋 Creating additional grid slides for remaining cards...")
    
    remaining_cards = len(card_images)
    slides_needed = math.ceil(remaining_cards / cards_per_slide)
    
    for slide_idx in range(slides_needed):
        # Add new slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Calculate grid layout
        grid_cols = 5
        grid_rows = 4
        margin = Inches(0.5)
        slide_width = Inches(10)  # Standard slide width
        slide_height = Inches(7.5)  # Standard slide height
        
        card_width = (slide_width - margin * 2) / grid_cols - Inches(0.1)
        card_height = card_width * 1.4  # Magic card aspect ratio
        
        cards_this_slide = min(cards_per_slide, remaining_cards)
        
        print(f"   📄 Creating grid slide {slide_idx + 1} with {cards_this_slide} cards")
        
        for card_in_slide in range(cards_this_slide):
            card_index = slide_idx * cards_per_slide + card_in_slide
            
            if card_index >= len(card_images):
                break
                
            # Calculate position in grid
            row = card_in_slide // grid_cols
            col = card_in_slide % grid_cols
            
            left = margin + col * (card_width + Inches(0.1))
            top = margin + row * (card_height + Inches(0.1))
            
            try:
                card_image = card_images[card_index]
                
                # Resize image
                resized_image = resize_image_for_slot(
                    card_image,
                    card_width / 914400,  # Convert to inches
                    card_height / 914400
                )
                
                # Add to slide
                picture = slide.shapes.add_picture(resized_image, left, top, card_width, card_height)
                
            except Exception as e:
                print(f"   ❌ Failed to add card {card_index + 1}: {e}")
        
        remaining_cards -= cards_this_slide

def main():
    print("🎴 Magic Cards PowerPoint Automation")
    print("=" * 50)
    
    # File paths
    template_file = "magic_cards_template.pptx"
    output_file = "magic_cards_completed.pptx"
    images_dir = "images"
    
    # Check if files exist
    if not os.path.exists(template_file):
        print(f"❌ Template file not found: {template_file}")
        return
    
    if not os.path.exists(images_dir):
        print(f"❌ Images directory not found: {images_dir}")
        return
    
    try:
        # Step 1: Analyze template
        prs, slide_info, total_slots = analyze_presentation_template(template_file)
        
        # Step 2: Get card images
        card_images = get_card_images(images_dir)
        
        if not card_images:
            print("❌ No card images found!")
            return
        
        # Step 3: Place cards in existing slots
        placed_cards = place_cards_in_template(prs, slide_info, card_images)
        
        # Step 4: Create additional slides if needed
        remaining_cards = card_images[placed_cards:]
        if remaining_cards:
            print(f"\n📊 {len(remaining_cards)} cards remaining - creating grid slides...")
            create_card_grid_slides(prs, remaining_cards)
        
        # Step 5: Save completed presentation
        prs.save(output_file)
        print(f"\n🎉 Completed presentation saved as: {output_file}")
        print(f"✅ Total cards placed: {len(card_images)}")
        print(f"✅ Template slots used: {min(placed_cards, total_slots)}")
        
        # Clean up temp files
        if os.path.exists("temp_resized"):
            import shutil
            shutil.rmtree("temp_resized")
            print("🧹 Cleaned up temporary files")
        
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()