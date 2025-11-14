#!/usr/bin/env python3
"""
Use the template structure for ALL cards by duplicating the template pattern
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image
from pathlib import Path

def get_template_slot_positions(template_slide):
    """Extract the positions and sizes of card slots from template"""
    card_slots = []
    
    for shape in template_slide.shapes:
        # Look for rectangular shapes that are card slots
        if (shape.shape_type in [1, 13] and  # Rectangle or Picture
            shape.width.inches > 1.0 and 
            shape.height.inches > 1.0):
            
            card_slots.append({
                'left': shape.left.inches,
                'top': shape.top.inches,
                'width': shape.width.inches,
                'height': shape.height.inches
            })
    
    # Sort slots by position (top to bottom, left to right)
    card_slots.sort(key=lambda x: (x['top'], x['left']))
    
    return card_slots

def place_card_in_slot(slide, card_path, slot_info):
    """Place a card image in the specified slot position"""
    try:
        # Load and resize the card image
        with Image.open(card_path) as img:
            # Calculate scaling to fit slot while maintaining aspect ratio
            slot_aspect = slot_info['width'] / slot_info['height']
            img_aspect = img.width / img.height
            
            if img_aspect > slot_aspect:
                # Image is wider, fit to width
                new_width = slot_info['width']
                new_height = slot_info['width'] / img_aspect
            else:
                # Image is taller, fit to height
                new_height = slot_info['height']
                new_width = slot_info['height'] * img_aspect
            
            # Center the image in the slot
            left = slot_info['left'] + (slot_info['width'] - new_width) / 2
            top = slot_info['top'] + (slot_info['height'] - new_height) / 2
            
            # Resize image for better quality
            target_width = int(new_width * 150)
            target_height = int(new_height * 150)
            
            resized_img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
            
            # Save temporary resized image
            temp_path = f"temp_card_{os.getpid()}.jpg"
            resized_img.save(temp_path, "JPEG", quality=90)
        
        # Add the card image to the slide
        slide.shapes.add_picture(
            temp_path,
            Inches(left),
            Inches(top),
            Inches(new_width),
            Inches(new_height)
        )
        
        # Clean up temporary file
        os.remove(temp_path)
        
        return True
        
    except Exception as e:
        print(f"  ❌ Failed to place card: {e}")
        return False

def create_presentation_with_template_pattern(template_file, output_file):
    """Create presentation using template pattern for ALL cards"""
    
    # Get all card images
    images_dir = "images"
    card_images = []
    
    if os.path.exists(images_dir):
        for filename in sorted(os.listdir(images_dir)):
            if filename.lower().endswith(('.jpg', '.jpeg', '.png')):
                card_images.append(os.path.join(images_dir, filename))
    
    total_cards = len(card_images)
    print(f"🎴 Found {total_cards} card images to place")
    
    # Load the template
    template_prs = Presentation(template_file)
    
    # Get the slot positions from the first slide
    if len(template_prs.slides) == 0:
        print("❌ Template has no slides!")
        return False
    
    template_slide = template_prs.slides[0]
    slot_positions = get_template_slot_positions(template_slide)
    
    print(f"📐 Found {len(slot_positions)} card slots in template")
    
    if len(slot_positions) == 0:
        print("❌ No card slots found in template!")
        return False
    
    cards_per_slide = len(slot_positions)
    slides_needed = math.ceil(total_cards / cards_per_slide)
    
    print(f"📄 Creating {slides_needed} slides with template pattern ({cards_per_slide} cards per slide)")
    
    # Create new presentation
    prs = Presentation()
    
    # Use blank slide layout
    blank_layout = prs.slide_layouts[5]
    
    # Process cards in groups
    for slide_num in range(slides_needed):
        # Add new slide
        slide = prs.slides.add_slide(blank_layout)
        
        # Add black background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
        
        # Calculate card indices for this slide
        start_idx = slide_num * cards_per_slide
        end_idx = min(start_idx + cards_per_slide, total_cards)
        
        print(f"\n📄 Slide {slide_num + 1}: Cards {start_idx + 1} to {end_idx}")
        
        # Place cards using template positions
        for i, card_idx in enumerate(range(start_idx, end_idx)):
            if i >= len(slot_positions):
                break
            
            card_path = card_images[card_idx]
            card_name = os.path.splitext(os.path.basename(card_path))[0]
            slot = slot_positions[i]
            
            if place_card_in_slot(slide, card_path, slot):
                print(f"  ✅ Placed {card_name}")
            else:
                print(f"  ❌ Failed to place {card_name}")
    
    # Save the presentation
    prs.save(output_file)
    print(f"\n💾 Saved presentation: {output_file}")
    
    return True

def main():
    """Create presentation using template for all cards"""
    
    template_file = "magic_cards_template.pptx"
    output_file = "Tifa_Lockhart_EDH_Deck_ALL_TEMPLATE.pptx"
    
    print(f"🎴 Creating Full Presentation Using Template Pattern")
    print("=" * 70)
    
    if not os.path.exists(template_file):
        print(f"❌ Template file not found: {template_file}")
        return False
    
    if not os.path.exists("images"):
        print(f"❌ Images directory not found")
        return False
    
    success = create_presentation_with_template_pattern(template_file, output_file)
    
    if success:
        print("\n" + "=" * 70)
        print(f"🎉 Presentation created successfully!")
        print(f"📁 File: {output_file}")
        print(f"🎴 All cards placed using template slot pattern!")
    else:
        print("❌ Failed to create presentation")
    
    return success

if __name__ == "__main__":
    main()