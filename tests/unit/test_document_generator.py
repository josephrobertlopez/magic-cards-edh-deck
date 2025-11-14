"""
Unit tests for magic_cards/document_generator.py

Tests cover:
- Content slot detection in PPTX templates
- Content placement with aspect-ratio preservation
- Template validation
"""

import pytest
import os
import tempfile
from unittest.mock import Mock, patch, MagicMock
from pptx import Presentation
from pptx.util import Inches
from magic_cards.document_generator import detect_content_slots, place_content_in_slot


# --- Fixtures ---

@pytest.fixture
def mock_template():
    """Create a mock PPTX template with two content slots"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add a slide with two rectangles (content slots)
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Vertical slot (2.5" × 3.5", Magic card aspect ratio)
    vertical_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1), Inches(1),
        Inches(2.5), Inches(3.5)
    )

    # Horizontal slot (3.5" × 2.5")
    horizontal_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(5), Inches(2),
        Inches(3.5), Inches(2.5)
    )

    # Save to temp file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as f:
        prs.save(f.name)
        return f.name


@pytest.fixture
def mock_empty_template():
    """Create a mock PPTX template with no content slots"""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as f:
        prs.save(f.name)
        return f.name


@pytest.fixture
def mock_image_path():
    """Create a small test image file"""
    from PIL import Image
    img = Image.new('RGB', (100, 140), color='red')
    with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as f:
        img.save(f.name)
        return f.name


# --- Unit Tests ---

def test_detect_content_slots_success(mock_template):
    """Test slot detection finds all valid content slots"""
    slots = detect_content_slots(mock_template)

    assert len(slots) == 2
    assert slots[0]['orientation'] == 'vertical'
    assert slots[1]['orientation'] == 'horizontal'

    # Clean up
    os.unlink(mock_template)


def test_detect_content_slots_sorted_by_position(mock_template):
    """Test slots are sorted top-to-bottom, left-to-right"""
    slots = detect_content_slots(mock_template)

    # First slot should be at (1, 1), second at (5, 2)
    assert slots[0]['left'] < slots[1]['left']
    assert slots[0]['top'] < slots[1]['top']

    os.unlink(mock_template)


def test_detect_content_slots_empty_template(mock_empty_template):
    """Test slot detection with template containing no valid slots"""
    slots = detect_content_slots(mock_empty_template)

    assert len(slots) == 0

    os.unlink(mock_empty_template)


def test_detect_content_slots_no_slides():
    """Test slot detection raises error for template with no slides"""
    prs = Presentation()
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as f:
        prs.save(f.name)
        template_path = f.name

    with pytest.raises(ValueError, match="Template has no slides"):
        detect_content_slots(template_path)

    os.unlink(template_path)


def test_place_content_in_slot(mock_template, mock_image_path):
    """Test content placement in slot with aspect-ratio preservation"""
    prs = Presentation(mock_template)
    slide = prs.slides[0]
    slots = detect_content_slots(mock_template)

    # Place image in first slot
    place_content_in_slot(slide, mock_image_path, slots[0])

    # Verify image was added (slide should have one more shape than slots)
    assert len(slide.shapes) > len(slots)

    # Clean up
    os.unlink(mock_template)
    os.unlink(mock_image_path)
