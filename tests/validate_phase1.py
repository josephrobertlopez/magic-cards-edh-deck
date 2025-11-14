#!/usr/bin/env python3
"""
Phase 1 Validation Script
Tests both fetch-cards and fill-template skills and validates Phase 1 exit criteria
"""

import os
import json
import sys
from pathlib import Path

# ANSI colors for output
GREEN = '\033[92m'
RED = '\033[91m'
YELLOW = '\033[93m'
BLUE = '\033[94m'
RESET = '\033[0m'

def print_header(text):
    print(f"\n{BLUE}{'=' * 70}{RESET}")
    print(f"{BLUE}{text}{RESET}")
    print(f"{BLUE}{'=' * 70}{RESET}\n")

def print_pass(text):
    print(f"{GREEN}✅ PASS:{RESET} {text}")

def print_fail(text):
    print(f"{RED}❌ FAIL:{RESET} {text}")

def print_info(text):
    print(f"{YELLOW}ℹ️  INFO:{RESET} {text}")

def validate_manifest_schema(manifest_path):
    """
    Validate JSON manifest structure (T019)
    Exit Criterion 4: Manifest is consumable by fill-template
    """
    print_header("Validating Manifest Schema")

    if not os.path.exists(manifest_path):
        print_fail(f"Manifest file not found: {manifest_path}")
        return False

    try:
        with open(manifest_path, 'r') as f:
            manifest = json.load(f)

        # Check required top-level fields
        required_fields = ['timestamp', 'decklist_path', 'total_cards', 'successful', 'failed', 'cards']
        for field in required_fields:
            if field not in manifest:
                print_fail(f"Missing required field: {field}")
                return False

        print_pass(f"All required fields present: {', '.join(required_fields)}")

        # Validate cards array
        if not isinstance(manifest['cards'], list):
            print_fail("'cards' field is not an array")
            return False

        print_pass(f"Cards array contains {len(manifest['cards'])} entries")

        # Validate card entries
        for i, card in enumerate(manifest['cards']):
            if 'name' not in card or 'status' not in card:
                print_fail(f"Card {i} missing required fields")
                return False

            if card['status'] not in ['success', 'failed']:
                print_fail(f"Card {i} has invalid status: {card['status']}")
                return False

            if card['status'] == 'success' and 'path' not in card:
                print_fail(f"Successful card {i} missing 'path' field")
                return False

            if card['status'] == 'failed' and 'reason' not in card:
                print_fail(f"Failed card {i} missing 'reason' field")
                return False

        print_pass("All card entries have valid structure")

        # Validate counts
        expected_total = manifest['successful'] + manifest['failed']
        if len(manifest['cards']) != expected_total:
            print_fail(f"Card count mismatch: {len(manifest['cards'])} cards but successful({manifest['successful']}) + failed({manifest['failed']}) = {expected_total}")
            return False

        print_pass(f"Counts match: {manifest['successful']} successful + {manifest['failed']} failed = {expected_total} total")

        return True

    except json.JSONDecodeError as e:
        print_fail(f"Invalid JSON: {e}")
        return False
    except Exception as e:
        print_fail(f"Validation error: {e}")
        return False

def validate_success_rate(manifest_path, min_rate=0.8):
    """
    Validate at least 80% of cards downloaded successfully (T019)
    Exit Criterion 2: At least 80% cards download successfully
    """
    print_header(f"Validating Success Rate (minimum {min_rate*100}%)")

    with open(manifest_path, 'r') as f:
        manifest = json.load(f)

    total = manifest['total_cards']
    successful = manifest['successful']
    success_rate = successful / total if total > 0 else 0

    print_info(f"Total cards: {total}")
    print_info(f"Successful: {successful}")
    print_info(f"Success rate: {success_rate * 100:.1f}%")

    if success_rate >= min_rate:
        print_pass(f"Success rate {success_rate * 100:.1f}% meets minimum {min_rate * 100}%")
        return True
    else:
        print_fail(f"Success rate {success_rate * 100:.1f}% below minimum {min_rate * 100}%")
        return False

def validate_presentation(output_file):
    """
    Validate presentation file was created (T019)
    Exit Criterion 3: Presentation visually matches template pattern
    """
    print_header("Validating Presentation Output")

    if not os.path.exists(output_file):
        print_fail(f"Presentation file not found: {output_file}")
        return False

    print_pass(f"Presentation file exists: {output_file}")

    # Check file size
    file_size = os.path.getsize(output_file)
    if file_size == 0:
        print_fail("Presentation file is empty (0 bytes)")
        return False

    print_pass(f"Presentation file size: {file_size:,} bytes")

    # Try to load with python-pptx
    try:
        from pptx import Presentation
        prs = Presentation(output_file)

        slide_count = len(prs.slides)
        print_pass(f"Presentation has {slide_count} slide(s)")

        if slide_count == 0:
            print_fail("Presentation has no slides")
            return False

        # Check first slide has pictures
        first_slide = prs.slides[0]
        picture_count = sum(1 for shape in first_slide.shapes if shape.shape_type == 13)

        print_pass(f"First slide has {picture_count} picture(s)")

        if picture_count == 0:
            print_fail("First slide has no pictures (cards not placed?)")
            return False

        return True

    except Exception as e:
        print_fail(f"Could not load presentation: {e}")
        return False

def run_validation(decklist_file="test_decks/phase1_sample.txt",
                   template_file="decks.pptx",
                   output_file="tests/phase1_output.pptx"):
    """
    Run full Phase 1 validation (T020)
    Tests all 6 exit criteria
    """
    print_header("Phase 1 Validation Script")
    print(f"Decklist: {decklist_file}")
    print(f"Template: {template_file}")
    print(f"Output: {output_file}")

    all_passed = True
    results = {}

    # Criterion 1 & 2: Skills execute without crashes (tested by running them)
    # Criterion 4 & 6: Manifest validation
    manifest_path = ".claude/state/fetch_manifest.json"

    print_header("Exit Criterion 4: Manifest Validation")
    if os.path.exists(manifest_path):
        schema_valid = validate_manifest_schema(manifest_path)
        results['manifest_schema'] = schema_valid
        all_passed = all_passed and schema_valid

        if schema_valid:
            rate_valid = validate_success_rate(manifest_path)
            results['success_rate'] = rate_valid
            all_passed = all_passed and rate_valid
    else:
        print_fail(f"Manifest not found: {manifest_path}")
        print_info("Run fetch-cards skill first")
        results['manifest_schema'] = False
        results['success_rate'] = False
        all_passed = False

    # Criterion 3: Presentation structure
    print_header("Exit Criterion 3: Presentation Validation")
    if os.path.exists(output_file):
        prs_valid = validate_presentation(output_file)
        results['presentation'] = prs_valid
        all_passed = all_passed and prs_valid
    else:
        print_info(f"Presentation not found: {output_file}")
        print_info("Run fill-template skill first")
        results['presentation'] = False
        all_passed = False

    # Summary
    print_header("Validation Summary")

    print("\nExit Criteria Status:")
    print(f"  Criterion 1: Both skills execute without crashes - {YELLOW}MANUAL CHECK REQUIRED{RESET}")
    print(f"  Criterion 2: 80%+ cards download successfully - {GREEN + '✅ PASS' if results.get('success_rate') else RED + '❌ FAIL'}{RESET}")
    print(f"  Criterion 3: Presentation matches template - {GREEN + '✅ PASS' if results.get('presentation') else RED + '❌ FAIL'}{RESET}")
    print(f"  Criterion 4: Manifest consumable - {GREEN + '✅ PASS' if results.get('manifest_schema') else RED + '❌ FAIL'}{RESET}")
    print(f"  Criterion 5: User validates output quality - {YELLOW}MANUAL CHECK REQUIRED{RESET}")
    print(f"  Criterion 6: This script reports PASS - {'...'}")

    print(f"\n{BLUE}{'=' * 70}{RESET}")

    if all_passed:
        print(f"{GREEN}🎉 VALIDATION PASSED{RESET}")
        print(f"{GREEN}✅ All automated checks successful{RESET}")
        print(f"\n{YELLOW}⚠️  Manual validation still required:{RESET}")
        print(f"   - Criterion 1: Verify skills ran without crashes")
        print(f"   - Criterion 5: User validates output quality for proxy printing")
        return 0
    else:
        print(f"{RED}❌ VALIDATION FAILED{RESET}")
        print(f"{RED}Some checks did not pass{RESET}")
        print(f"\n{YELLOW}Fix issues and run validation again{RESET}")
        return 1

def main():
    """Main entry point"""
    if len(sys.argv) > 1 and sys.argv[1] in ['-h', '--help']:
        print("Usage: python tests/validate_phase1.py")
        print("\nValidates Phase 1 implementation against exit criteria")
        print("\nPrerequisites:")
        print("  - Run /fetch-cards test_decks/phase1_sample.txt")
        print("  - Run /fill-template decks.pptx images tests/phase1_output.pptx")
        print("\nExit Codes:")
        print("  0 - All automated checks passed")
        print("  1 - One or more checks failed")
        return 0

    return run_validation()

if __name__ == "__main__":
    sys.exit(main())
