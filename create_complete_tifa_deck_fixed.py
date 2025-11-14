#!/usr/bin/env python3
"""
Create complete Tifa Lockhart EDH deck presentation with all 100 cards
Uses the PNG template layout (2 vertical + 6 horizontal slots per slide)
"""

import os
import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def get_card_images():
    """Get all card images from the images folder"""
    images_dir = Path("images")
    if not images_dir.exists():
        raise FileNotFoundError("Images directory not found!")
    
    # Get all .jpg files
    card_files = list(images_dir.glob("*.jpg"))
    card_files.sort()  # Sort alphabetically
    
    print(f"📁 Found {len(card_files)} card images")
    return card_files

def resize_image_for_slot(image_path, target_width, target_height, orientation='horizontal'):
    """
    Resize image to fit the slot while maintaining aspect ratio
    """
    try:
        with Image.open(image_path) as img:
            # Magic card aspect ratio is approximately 2.5:3.5 (width:height)
            card_aspect_ratio = 2.5 / 3.5  # ~0.714
            
            if orientation == 'vertical':
                # For vertical slots, use the card in portrait orientation
                slot_aspect_ratio = target_width / target_height
            else:
                # For horizontal slots, rotate the card to landscape
                slot_aspect_ratio = target_width / target_height
                # Swap dimensions for landscape orientation of card
                card_aspect_ratio = 3.5 / 2.5  # ~1.4
            
            # Calculate the best fit
            if card_aspect_ratio > slot_aspect_ratio:
                # Card is wider, fit to width
                new_width = target_width
                new_height = target_width / card_aspect_ratio
            else:
                # Card is taller, fit to height
                new_height = target_height
                new_width = target_height * card_aspect_ratio
            
            return new_width, new_height
            
    except Exception as e:
        print(f"  ⚠️ Error resizing {image_path}: {e}")
        return target_width, target_height

def create_complete_tifa_deck_presentation(template_file, card_images, output_file):
    """Create complete Tifa Lockhart EDH deck presentation with 100 cards"""
    
    try:
        # Load the template
        prs = Presentation(template_file)
        print(f"📋 Loaded template with {len(prs.slides)} slides")
        
        total_cards = len(card_images)
        cards_per_slide = 8  # 2 vertical + 6 horizontal slots
        slides_needed = math.ceil(total_cards / cards_per_slide)
        
        print(f"🎴 Processing {total_cards} cards")
        print(f"📄 Need {slides_needed} slides ({cards_per_slide} cards per slide)")
        
        # Slot positions based on PNG template analysis
        # 2 vertical slots on the left, 6 horizontal slots on the right in 2x3 grid
        slot_positions = [
            # 2 vertical slots (left column) - portrait orientation
            {'x': 0.5, 'y': 1.0, 'width': 2.5, 'height': 3.5, 'orientation': 'vertical'},
            {'x': 0.5, 'y': 5.0, 'width': 2.5, 'height': 3.5, 'orientation': 'vertical'},
            
            # 6 horizontal slots (right side, 2x3 grid) - landscape orientation  
            {'x': 4.0, 'y': 0.5, 'width': 2.5, 'height': 1.8, 'orientation': 'horizontal'},
            {'x': 7.0, 'y': 0.5, 'width': 2.5, 'height': 1.8, 'orientation': 'horizontal'},
            {'x': 4.0, 'y': 2.8, 'width': 2.5, 'height': 1.8, 'orientation': 'horizontal'},
            {'x': 7.0, 'y': 2.8, 'width': 2.5, 'height': 1.8, 'orientation': 'horizontal'},
            {'x': 4.0, 'y': 5.1, 'width': 2.5, 'height': 1.8, 'orientation': 'horizontal'},
            {'x': 7.0, 'y': 5.1, 'width': 2.5, 'height': 1.8, 'orientation': 'horizontal'},
        ]
        
        card_index = 0
        
        for slide_num in range(slides_needed):
            print(f"\n📄 Creating slide {slide_num + 1}/{slides_needed}")
            
            if slide_num == 0:
                # Use the first slide (modify existing)
                slide = prs.slides[0]
                print("  📝 Using existing slide 1")
            else:
                # Duplicate the first slide layout for additional slides
                slide_layout = prs.slide_layouts[0]  # Use the first layout
                slide = prs.slides.add_slide(slide_layout)
                print("  ➕ Added new slide")
            
            # Place cards in the 8 slots for this slide
            for slot_index, slot in enumerate(slot_positions):
                if card_index >= total_cards:
                    break
                
                card_file = card_images[card_index]
                print(f"  🎴 Adding card {card_index + 1}: {card_file.name} to slot {slot_index + 1}")
                
                try:
                    # Calculate optimal size for this slot
                    slot_width = Inches(slot['width'])
                    slot_height = Inches(slot['height'])
                    
                    # Get optimal dimensions while maintaining aspect ratio
                    optimal_width, optimal_height = resize_image_for_slot(
                        card_file, 
                        slot['width'], 
                        slot['height'],
                        slot['orientation']
                    )
                    
                    # Add the image
                    pic = slide.shapes.add_picture(
                        str(card_file),
                        Inches(slot['x']),
                        Inches(slot['y']),
                        Inches(optimal_width),
                        Inches(optimal_height)
                    )
                    
                    print(f"    ✅ Placed at ({slot['x']:.1f}, {slot['y']:.1f}) - {optimal_width:.1f}x{optimal_height:.1f} inches ({slot['orientation']})")
                    
                except Exception as e:
                    print(f"    ❌ Error adding {card_file.name}: {e}")
                
                card_index += 1
        
        # Save the presentation
        prs.save(output_file)
        print(f"\n💾 Saved presentation: {output_file}")
        print(f"📊 Final stats: {slides_needed} slides, {card_index} cards placed")
        
        return True
        
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        return False

def main():
    """Main function"""
    template_file = "magic_cards_template.pptx" 
    output_file = "Tifa_Lockhart_Complete_EDH_Deck.pptx"
    
    print("🎯 Creating Complete Tifa Lockhart EDH Deck Presentation")
    print("=" * 70)
    
    try:
        # Get all card images
        card_images = get_card_images()
        
        if len(card_images) == 0:
            print("❌ No card images found!")
            return
        
        # Create the presentation
        success = create_complete_tifa_deck_presentation(template_file, card_images, output_file)
        
        if success:
            print(f"\n🎉 Successfully created {output_file}")
            print(f"🎴 Complete EDH deck with {len(card_images)} cards!")
        else:
            print("\n❌ Failed to create presentation")
            
    except Exception as e:
        print(f"❌ Error: {e}")

if __name__ == "__main__":
    main()