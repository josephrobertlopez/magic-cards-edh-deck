#!/usr/bin/env python3
"""
Magic Cards Template Automation - DARK SLOTS ONLY VERSION
Only fills the 8 dark slots on Slide 1, leaves everything else untouched.
No additional slides created.
"""

import os
import glob
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

def resize_card_for_slot(card_path, slot_width_inches, slot_height_inches):
    """
    Resize a Magic card image to fit a slot while maintaining aspect ratio
    """
    try:
        with Image.open(card_path) as img:
            # Standard Magic card aspect ratio is 2.5:3.5 (width:height) = 0.714
            card_aspect = 2.5 / 3.5
            
            # Calculate size to fit the slot
            slot_aspect = slot_width_inches / slot_height_inches
            
            if card_aspect > slot_aspect:
                # Image is wider than slot - fit to width
                new_width = slot_width_inches
                new_height = slot_width_inches / card_aspect
            else:
                # Image is taller than slot - fit to height
                new_height = slot_height_inches
                new_width = slot_height_inches * card_aspect
            
            # Convert to pixels (assuming 96 DPI)
            new_width_px = int(new_width * 96)
            new_height_px = int(new_height * 96)
            
            # Resize the image
            img_resized = img.resize((new_width_px, new_height_px), Image.Resampling.LANCZOS)
            
            # Save to bytes
            img_stream = io.BytesIO()
            img_resized.save(img_stream, format='PNG')
            img_stream.seek(0)
            
            return img_stream, new_width, new_height
            
    except Exception as e:
        print(f"    ❌ Error resizing {card_path}: {e}")
        return None, 0, 0

def place_cards_in_dark_slots_only(template_file, card_images, output_file):
    """
    Place cards ONLY in the 8 dark slots on Slide 1.
    Leave Slide 2 and everything else completely untouched.
    """
    print(f"🔧 Opening template: {template_file}")
    prs = Presentation(template_file)
    
    if len(prs.slides) < 1:
        print("❌ Template has no slides!")
        return
    
    # ONLY process Slide 1 - the one with 8 dark slots
    slide = prs.slides[0]  # First slide (index 0)
    print(f"\n📋 Processing Slide 1 ONLY ({len(slide.shapes)} shapes)")
    
    # Find all card slots on Slide 1
    card_slots = []
    for shape_num, shape in enumerate(slide.shapes):
        width_inches = shape.width / 914400
        height_inches = shape.height / 914400
        aspect_ratio = width_inches / height_inches if height_inches > 0 else 0
        
        # Detect Magic card slots by size and aspect ratio (vertical orientation)
        if 0.65 < aspect_ratio < 0.8 and width_inches > 2.0 and height_inches > 3.0:
            card_slots.append({
                'shape': shape,
                'shape_num': shape_num + 1,
                'width_inches': width_inches,
                'height_inches': height_inches
            })
    
    print(f"  🎴 Found {len(card_slots)} dark card slots on Slide 1")
    
    if len(card_slots) != 8:
        print(f"  ⚠️  Expected 8 dark slots, found {len(card_slots)}")
    
    # Use only the first 8 cards to fill the 8 dark slots
    cards_to_use = card_images[:8]
    
    # Replace each dark slot with a card image
    for i, slot_info in enumerate(card_slots):
        if i >= len(cards_to_use):
            print(f"  ⚠️  Ran out of cards at slot {i + 1}")
            break
            
        card_path = cards_to_use[i]
        card_name = os.path.basename(card_path)
        
        print(f"  🎴 Dark Slot {i + 1}: {card_name}")
        
        # Resize card for this specific slot
        img_stream, img_width, img_height = resize_card_for_slot(
            card_path, 
            slot_info['width_inches'], 
            slot_info['height_inches']
        )
        
        if img_stream:
            try:
                # Remove the original dark slot shape
                shape_element = slot_info['shape'].element
                shape_element.getparent().remove(shape_element)
                
                # Add the card image in the same position
                left = slot_info['shape'].left
                top = slot_info['shape'].top
                
                # Add picture with proper size
                new_shape = slide.shapes.add_picture(
                    img_stream, 
                    left, 
                    top, 
                    Inches(img_width),
                    Inches(img_height)
                )
                
                print(f"    ✅ Placed {card_name} in dark slot {i + 1}")
                
            except Exception as e:
                print(f"    ❌ Error placing {card_name}: {e}")
    
    # Save the presentation (NO additional slides added)
    print(f"\n💾 Saving to: {output_file}")
    prs.save(output_file)
    
    print(f"\n🎉 Dark slots processing complete!")
    print(f"  📊 Total slides: {len(prs.slides)} (unchanged)")
    print(f"  🎴 Dark slots filled: {min(len(card_slots), len(cards_to_use))}")
    print(f"  📄 Slide 2: Left completely untouched")
    print(f"  🎯 Cards used: {len(cards_to_use)} of {len(card_images)} total")

def main():
    # File paths
    template_file = "magic_cards_template.pptx"
    output_file = "magic_cards_completed_DARK_SLOTS_ONLY.pptx"
    images_dir = "images"
    
    # Check if template exists
    if not os.path.exists(template_file):
        print(f"❌ Template file not found: {template_file}")
        return
    
    # Check if images directory exists
    if not os.path.exists(images_dir):
        print(f"❌ Images directory not found: {images_dir}")
        return
    
    # Get all card images (sorted by filename)
    card_images = sorted(glob.glob(os.path.join(images_dir, "*.jpg")))
    if not card_images:
        card_images = sorted(glob.glob(os.path.join(images_dir, "*.png")))
    
    if not card_images:
        print(f"❌ No card images found in {images_dir}")
        return
    
    print(f"🎴 Found {len(card_images)} card images")
    print(f"📁 Template: {template_file}")
    print(f"💾 Output: {output_file}")
    print(f"🎯 Processing ONLY the 8 dark slots...")
    
    # Process ONLY the dark slots
    place_cards_in_dark_slots_only(template_file, card_images, output_file)

if __name__ == "__main__":
    main()