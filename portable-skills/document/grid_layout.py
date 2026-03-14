#!/usr/bin/env python3
"""Generate print-ready grid layouts from images. Configurable rows/cols, margins, and page size."""

import argparse
import glob
import math
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu


# Standard page sizes in inches (landscape)
PAGE_SIZES = {
    "letter": (11, 8.5),
    "a4": (11.69, 8.27),
    "legal": (14, 8.5),
    "tabloid": (17, 11),
}


def create_grid_pptx(image_paths, output_path, cols=4, rows=2, page_size="letter",
                     margin=0.25, spacing=0.05, bg_color="white"):
    """Generate PPTX with images arranged in a grid layout.

    Args:
        image_paths: List of image file paths
        output_path: Output PPTX path
        cols: Columns per page
        rows: Rows per page
        page_size: Page size name or (width, height) tuple
        margin: Page margin in inches
        spacing: Gap between images in inches
        bg_color: Background color
    """
    if isinstance(page_size, str):
        width, height = PAGE_SIZES.get(page_size.lower(), PAGE_SIZES["letter"])
    else:
        width, height = page_size

    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    cards_per_page = cols * rows
    available_w = width - (2 * margin) - ((cols - 1) * spacing)
    available_h = height - (2 * margin) - ((rows - 1) * spacing)
    cell_w = available_w / cols
    cell_h = available_h / rows

    total_pages = math.ceil(len(image_paths) / cards_per_page)

    for page in range(total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        start = page * cards_per_page
        page_images = image_paths[start:start + cards_per_page]

        for idx, img_path in enumerate(page_images):
            row = idx // cols
            col = idx % cols

            left = margin + col * (cell_w + spacing)
            top = margin + row * (cell_h + spacing)

            try:
                # Get image aspect ratio for proper fitting
                with Image.open(img_path) as img:
                    img_w, img_h = img.size
                    img_aspect = img_w / img_h
                    cell_aspect = cell_w / cell_h

                    if img_aspect > cell_aspect:
                        # Image wider than cell — fit to width
                        fit_w = cell_w
                        fit_h = cell_w / img_aspect
                        top_offset = (cell_h - fit_h) / 2
                        left_offset = 0
                    else:
                        # Image taller than cell — fit to height
                        fit_h = cell_h
                        fit_w = cell_h * img_aspect
                        left_offset = (cell_w - fit_w) / 2
                        top_offset = 0

                slide.shapes.add_picture(
                    img_path,
                    Inches(left + left_offset),
                    Inches(top + top_offset),
                    Inches(fit_w),
                    Inches(fit_h),
                )
            except Exception as e:
                print(f"Warning: Skipped {img_path}: {e}", file=sys.stderr)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    print(f"Generated: {output_path}")
    print(f"  {len(image_paths)} images across {total_pages} pages ({cols}x{rows} grid)")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate grid layouts from images")
    parser.add_argument("images", nargs="+", help="Image files or glob patterns")
    parser.add_argument("--output", "-o", default="grid_output.pptx", help="Output PPTX path")
    parser.add_argument("--cols", type=int, default=4, help="Columns per page (default: 4)")
    parser.add_argument("--rows", type=int, default=2, help="Rows per page (default: 2)")
    parser.add_argument("--page-size", default="letter", choices=list(PAGE_SIZES.keys()),
                        help="Page size (default: letter)")
    parser.add_argument("--margin", type=float, default=0.25, help="Page margin in inches")
    parser.add_argument("--spacing", type=float, default=0.05, help="Gap between images in inches")
    args = parser.parse_args()

    # Expand glob patterns
    image_paths = []
    for pattern in args.images:
        expanded = glob.glob(pattern)
        if expanded:
            image_paths.extend(sorted(expanded))
        elif os.path.exists(pattern):
            image_paths.append(pattern)
        else:
            print(f"Warning: No match for {pattern}", file=sys.stderr)

    if not image_paths:
        print("Error: No images found", file=sys.stderr)
        sys.exit(1)

    create_grid_pptx(
        image_paths, args.output,
        cols=args.cols, rows=args.rows,
        page_size=args.page_size,
        margin=args.margin, spacing=args.spacing,
    )


if __name__ == "__main__":
    main()
