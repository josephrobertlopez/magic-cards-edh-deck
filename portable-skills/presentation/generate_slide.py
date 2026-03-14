#!/usr/bin/env python3
"""Generate PPTX presentations from image folders. Configurable grid, sizing, and templates."""

import argparse
import glob
import json
import math
import os
import sys

from pptx import Presentation
from pptx.util import Inches


def generate_slides(image_paths, output_path, per_slide=9, template=None,
                    slide_width=10.0, slide_height=7.5):
    """Generate PPTX with images in grid layout.

    Args:
        image_paths: List of image file paths
        output_path: Output PPTX path
        per_slide: Images per slide (default: 9 = 3x3)
        template: Optional template PPTX path
        slide_width: Slide width in inches
        slide_height: Slide height in inches

    Returns:
        Result dict
    """
    if template and os.path.exists(template):
        prs = Presentation(template)
    else:
        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)

    grid_size = int(math.sqrt(per_slide))
    if grid_size * grid_size < per_slide:
        grid_cols = grid_size + 1
        grid_rows = math.ceil(per_slide / grid_cols)
    else:
        grid_cols = grid_size
        grid_rows = grid_size

    margin = 0.25
    available_w = slide_width - (2 * margin)
    available_h = slide_height - (2 * margin)
    cell_w = available_w / grid_cols
    cell_h = available_h / grid_rows

    total_slides = math.ceil(len(image_paths) / per_slide)

    for slide_idx in range(total_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        start = slide_idx * per_slide
        slide_images = image_paths[start:start + per_slide]

        for idx, img_path in enumerate(slide_images):
            row = idx // grid_cols
            col = idx % grid_cols
            left = Inches(margin + col * cell_w)
            top = Inches(margin + row * cell_h)

            try:
                slide.shapes.add_picture(img_path, left, top,
                                         Inches(cell_w), Inches(cell_h))
            except Exception as e:
                print(f"Warning: Skipped {img_path}: {e}", file=sys.stderr)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    result = {
        "pptx_path": output_path,
        "total_images": len(image_paths),
        "total_slides": total_slides,
        "cards_per_slide": per_slide,
        "grid": f"{grid_cols}x{grid_rows}",
        "status": "success",
    }
    return result


def main():
    parser = argparse.ArgumentParser(description="Generate PPTX from images")
    parser.add_argument("images", nargs="+", help="Image files, directories, or glob patterns")
    parser.add_argument("--output", "-o", default="slides.pptx", help="Output PPTX path")
    parser.add_argument("--per-slide", type=int, default=9, help="Images per slide (default: 9)")
    parser.add_argument("--template", help="Template PPTX file")
    parser.add_argument("--width", type=float, default=10.0, help="Slide width in inches")
    parser.add_argument("--height", type=float, default=7.5, help="Slide height in inches")
    args = parser.parse_args()

    # Collect image paths
    image_paths = []
    image_exts = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff"}
    for item in args.images:
        if os.path.isdir(item):
            for f in sorted(os.listdir(item)):
                if os.path.splitext(f)[1].lower() in image_exts:
                    image_paths.append(os.path.join(item, f))
        else:
            expanded = glob.glob(item)
            image_paths.extend(sorted(expanded) if expanded else [item])

    if not image_paths:
        print("Error: No images found", file=sys.stderr)
        sys.exit(1)

    result = generate_slides(
        image_paths, args.output,
        per_slide=args.per_slide,
        template=args.template,
        slide_width=args.width,
        slide_height=args.height,
    )

    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
