#!/usr/bin/env python3
"""
FIXED Magic Cards Template Automation
=====================================
This script PROPERLY follows the original template structure:
- Slide 1: 8 card slots (shapes 1-8)
- Slide 2: 16 card slots (shapes 1-16)
- Total: 24 template slots

Remaining cards (158-24=134) go to additional grid slides.
"""

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import glob

def get_card_images():
    """Get all card images from the images folder"""
    image_files = glob.glob("images/*.jpg")
    image_files.sort()  # Sort to maintain consistent order
    print(f"🃏 Found {len(image_files)} card images")
    return image_files

def place_cards_in_template_properly(template_file, card_images, output_file):
    """Place cards in the EXACT template structure without changing anything"""
    print(f"🔧 Opening template: {template_file}")
    prs = Presentation(template_file)
    
    card_index = 0
    
    # Process slides that already exist in template
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n📋 Processing Slide {slide_num} ({len(slide.shapes)} shapes)")
        
        for shape_num, shape in enumerate(slide.shapes, 1):
            if card_index >= len(card_images):
                break
                
            # Check if this shape looks like a card slot
            width_inches = shape.width / 914400
            height_inches = shape.height / 914400
            aspect_ratio = width_inches / height_inches if height_inches > 0 else 0
            
            # Magic card aspect ratio is around 0.71 (2.5/3.5)
            if 0.65 < aspect_ratio < 0.8 and width_inches > 2.0 and height_inches > 3.0:
                print(f"   🎴 Filling slot {shape_num}: {width_inches:.2f}\" x {height_inches:.2f}\"")
                
                # Place the card image
                card_image_path = card_images[card_index]
                
                try:
                    # For AUTO_SHAPE, replace with picture
                    if shape.shape_type == 1:  # AUTO_SHAPE
                        # Get position and size
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        # Remove the shape
                        shape._element.getparent().remove(shape._element)
                        
                        # Add picture in same position
                        slide.shapes.add_picture(
                            card_image_path, left, top, width, height
                        )
                    
                    # For PICTURE, replace the image
                    elif shape.shape_type == 13:  # PICTURE
                        # Get position and size  
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        # Remove the old picture
                        shape._element.getparent().remove(shape._element)
                        
                        # Add new picture in same position
                        slide.shapes.add_picture(
                            card_image_path, left, top, width, height
                        )
                    
                    print(f"   ✅ Placed: {os.path.basename(card_image_path)}")
                    card_index += 1
                    
                except Exception as e:
                    print(f"   ❌ Failed to place {card_image_path}: {e}")
            else:
                print(f"   ⏭️ Skipping shape {shape_num}: {width_inches:.2f}\" x {height_inches:.2f}\" (not card-sized)")
    
    print(f"\n📊 Template slots used: {card_index}/24")
    
    # Add remaining cards to new grid slides
    remaining_cards = card_images[card_index:]
    if remaining_cards:
        print(f"📄 Creating additional slides for {len(remaining_cards)} remaining cards")
        create_card_grid_slides(prs, remaining_cards)
    
    # Save the result
    print(f"💾 Saving to: {output_file}")
    prs.save(output_file)
    
    return {
        "template_slots_used": card_index,
        "additional_cards": len(remaining_cards),
        "total_slides": len(prs.slides)
    }

def create_card_grid_slides(prs, remaining_cards, cards_per_slide=20):
    """Create grid slides for overflow cards"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    
    cards_placed = 0
    for i in range(0, len(remaining_cards), cards_per_slide):
        slide_cards = remaining_cards[i:i + cards_per_slide]
        
        print(f"📄 Creating grid slide with {len(slide_cards)} cards")
        slide = prs.slides.add_slide(slide_layout)
        
        # Create grid layout (4 columns x 5 rows)
        cols = 4
        rows = 5
        
        # Slide dimensions (assuming standard slide)
        slide_width = Inches(10)
        slide_height = Inches(7.5)
        
        # Card dimensions for grid
        card_width = Inches(2.0)
        card_height = Inches(2.8)
        
        # Spacing
        margin_x = (slide_width - (cols * card_width)) / (cols + 1)
        margin_y = (slide_height - (rows * card_height)) / (rows + 1)
        
        for j, card_path in enumerate(slide_cards):
            row = j // cols
            col = j % cols
            
            left = margin_x + col * (card_width + margin_x)
            top = margin_y + row * (card_height + margin_y)
            
            try:
                slide.shapes.add_picture(card_path, left, top, card_width, card_height)
                cards_placed += 1
            except Exception as e:
                print(f"❌ Failed to add {card_path}: {e}")
    
    print(f"✅ Added {cards_placed} cards to grid slides")

def main():
    print("🎴 FIXED Magic Cards Template Automation")
    print("=" * 50)
    
    # Get card images
    card_images = get_card_images()
    
    if not card_images:
        print("❌ No card images found in images/ folder!")
        return
    
    # Process the template properly
    result = place_cards_in_template_properly(
        "magic_cards_template.pptx",
        card_images,
        "magic_cards_completed_FIXED.pptx"
    )
    
    print("\n🎉 FIXED AUTOMATION COMPLETE!")
    print(f"✅ Template slots used: {result['template_slots_used']}")
    print(f"✅ Additional cards in grid: {result['additional_cards']}")
    print(f"✅ Total slides: {result['total_slides']}")
    print(f"📁 Output: magic_cards_completed_FIXED.pptx")

if __name__ == "__main__":
    main()