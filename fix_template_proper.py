#!/usr/bin/env python3
"""
Fix the presentation to PROPERLY follow the original template structure
by working with the existing shapes in the PowerPoint template
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from pathlib import Path

def analyze_template_shapes(template_file):
    """Analyze the existing shapes in the template"""
    prs = Presentation(template_file)
    
    print(f"🔍 Analyzing template: {template_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 Slide {slide_idx + 1}:")
        print(f"   Total shapes: {len(slide.shapes)}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            print(f"   Shape {shape_idx + 1}:")
            print(f"     Type: {shape.shape_type}")
            print(f"     Position: ({shape.left.inches:.2f}, {shape.top.inches:.2f})")
            print(f"     Size: {shape.width.inches:.2f} x {shape.height.inches:.2f}")
            
            # Check if it's a placeholder or text box
            if hasattr(shape, 'text'):
                print(f"     Text: '{shape.text[:50]}'")
            
            # Check if it's a picture
            if shape.shape_type == 13:  # Picture
                print(f"     Picture shape detected")

def replace_template_shapes_with_cards(template_file, output_file):
    """Replace the existing shapes in the template with card images"""
    
    # Get all card images
    images_dir = "images"
    card_images = []
    
    if os.path.exists(images_dir):
        for filename in sorted(os.listdir(images_dir)):
            if filename.lower().endswith(('.jpg', '.jpeg', '.png')):
                card_images.append(os.path.join(images_dir, filename))
    
    total_cards = len(card_images)
    print(f"🎴 Found {total_cards} card images to place")
    
    # Load the template
    prs = Presentation(template_file)
    
    print(f"📊 Template has {len(prs.slides)} slides")
    
    # Work with slide 1 (the one with dark slots)
    if len(prs.slides) >= 1:
        slide = prs.slides[0]
        print(f"🖼️ Working with slide 1 - {len(slide.shapes)} shapes found")
        
        # Find rectangular shapes that look like card slots
        card_slots = []
        for shape in slide.shapes:
            # Look for shapes that could be card slots (rectangles)
            if (shape.shape_type in [1, 13] and  # Rectangle or Picture
                shape.width.inches > 1.0 and  # Reasonable card width
                shape.height.inches > 1.0):   # Reasonable card height
                
                card_slots.append({
                    'shape': shape,
                    'left': shape.left.inches,
                    'top': shape.top.inches,
                    'width': shape.width.inches,
                    'height': shape.height.inches
                })
        
        print(f"🎯 Found {len(card_slots)} potential card slots")
        
        # Sort slots by position (top to bottom, left to right)
        card_slots.sort(key=lambda x: (x['top'], x['left']))
        
        # Replace shapes with card images
        cards_placed = 0
        
        for i, slot in enumerate(card_slots):
            if cards_placed >= total_cards:
                break
                
            card_path = card_images[cards_placed]
            card_name = os.path.splitext(os.path.basename(card_path))[0]
            
            try:
                # Get the original shape for positioning reference
                original_shape = slot['shape']
                
                # Remove the original shape
                shape_element = original_shape.element
                shape_element.getparent().remove(shape_element)
                
                # Load and resize the card image
                with Image.open(card_path) as img:
                    # Calculate scaling to fit slot while maintaining aspect ratio
                    slot_aspect = slot['width'] / slot['height']
                    img_aspect = img.width / img.height
                    
                    if img_aspect > slot_aspect:
                        # Image is wider, fit to width
                        new_width = slot['width']
                        new_height = slot['width'] / img_aspect
                    else:
                        # Image is taller, fit to height
                        new_height = slot['height']
                        new_width = slot['height'] * img_aspect
                    
                    # Center the image in the slot
                    left = slot['left'] + (slot['width'] - new_width) / 2
                    top = slot['top'] + (slot['height'] - new_height) / 2
                    
                    # Resize image for better quality
                    target_width = int(new_width * 150)  # Higher DPI for better quality
                    target_height = int(new_height * 150)
                    
                    resized_img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
                    
                    # Save temporary resized image
                    temp_path = f"temp_slot_{i}.jpg"
                    resized_img.save(temp_path, "JPEG", quality=90)
                
                # Add the card image to the slide
                slide.shapes.add_picture(
                    temp_path,
                    Inches(left),
                    Inches(top),
                    Inches(new_width),
                    Inches(new_height)
                )
                
                # Clean up temporary file
                os.remove(temp_path)
                
                print(f"  ✅ Placed {card_name} in slot {i+1}")
                cards_placed += 1
                
            except Exception as e:
                print(f"  ❌ Failed to place {card_name}: {e}")
    
    # Handle overflow cards on additional slides
    remaining_cards = total_cards - cards_placed
    
    if remaining_cards > 0:
        print(f"📄 Creating additional slides for {remaining_cards} remaining cards")
        
        # Use slide 2 as template for additional slides if it exists
        if len(prs.slides) >= 2:
            base_layout = prs.slide_layouts[5]  # Blank layout
        else:
            base_layout = prs.slide_layouts[0]  # Title slide layout
        
        cards_per_additional_slide = 8  # Standard grid layout
        additional_slides_needed = math.ceil(remaining_cards / cards_per_additional_slide)
        
        for slide_num in range(additional_slides_needed):
            # Create new slide
            new_slide = prs.slides.add_slide(base_layout)
            
            # Create grid layout for overflow cards
            start_idx = cards_placed + (slide_num * cards_per_additional_slide)
            end_idx = min(start_idx + cards_per_additional_slide, total_cards)
            
            # 2x4 grid layout
            grid_positions = [
                {'x': 1.0, 'y': 1.0}, {'x': 6.0, 'y': 1.0},
                {'x': 1.0, 'y': 3.0}, {'x': 6.0, 'y': 3.0},
                {'x': 1.0, 'y': 5.0}, {'x': 6.0, 'y': 5.0},
                {'x': 1.0, 'y': 7.0}, {'x': 6.0, 'y': 7.0},
            ]
            
            for i, card_idx in enumerate(range(start_idx, end_idx)):
                if i >= len(grid_positions):
                    break
                    
                card_path = card_images[card_idx]
                card_name = os.path.splitext(os.path.basename(card_path))[0]
                
                try:
                    pos = grid_positions[i]
                    
                    # Load and resize image
                    with Image.open(card_path) as img:
                        # Standard card size for overflow
                        target_width = int(2.5 * 150)
                        target_height = int(3.5 * 150)
                        
                        resized_img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
                        
                        temp_path = f"temp_overflow_{slide_num}_{i}.jpg"
                        resized_img.save(temp_path, "JPEG", quality=90)
                    
                    # Add to slide
                    new_slide.shapes.add_picture(
                        temp_path,
                        Inches(pos['x']),
                        Inches(pos['y']),
                        Inches(2.5),
                        Inches(3.5)
                    )
                    
                    os.remove(temp_path)
                    print(f"  ✅ Added {card_name} to overflow slide {slide_num + 1}")
                    
                except Exception as e:
                    print(f"  ❌ Failed to add {card_name}: {e}")
    
    # Save the presentation
    prs.save(output_file)
    print(f"💾 Saved corrected presentation: {output_file}")
    
    return True

def main():
    """Fix the template to properly follow the original structure"""
    
    template_file = "magic_cards_template.pptx"
    output_file = "Tifa_Lockhart_EDH_Deck_TEMPLATE_FIXED.pptx"
    
    print(f"🎴 Fixing Template Structure for Tifa Lockhart EDH Deck")
    print("=" * 70)
    
    if not os.path.exists(template_file):
        print(f"❌ Template file not found: {template_file}")
        return False
    
    if not os.path.exists("images"):
        print(f"❌ Images directory not found")
        return False
    
    # First analyze the template
    print("🔍 Step 1: Analyzing template structure...")
    analyze_template_shapes(template_file)
    
    print("\n🔧 Step 2: Replacing template shapes with cards...")
    success = replace_template_shapes_with_cards(template_file, output_file)
    
    if success:
        print("\n" + "=" * 70)
        print(f"🎉 Template-corrected presentation created!")
        print(f"📁 File: {output_file}")
        print(f"🎴 Now properly follows the original template structure!")
    else:
        print("❌ Failed to create corrected presentation")

if __name__ == "__main__":
    main()