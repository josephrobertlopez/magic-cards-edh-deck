#!/usr/bin/env python3
"""
Grid Layout Generator

Create print-ready PPTX/PDF with images arranged in a configurable grid.
Supports any NxM layout on any paper size. Defaults to US Letter with
MTG card dimensions (2x4 landscape grid).

Works with any images: Scryfall cards, custom art, tokens, alters, etc.

Usage:
    # Default 2x4 MTG layout on Letter paper
    python3 grid_layout.py images/card1.jpg images/card2.jpg ...

    # From a directory of images
    python3 grid_layout.py images/*.jpg

    # Custom 3x3 grid
    python3 grid_layout.py --rows 3 --cols 3 images/*.png

    # Custom paper size (A4)
    python3 grid_layout.py --width 8.27 --height 11.69 images/*.jpg

    # No rotation (portrait cards)
    python3 grid_layout.py --no-rotate images/*.jpg
"""

import sys
import json
import argparse
import subprocess
import math
from pathlib import Path
from typing import Dict, Any, List, Optional

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image


# Preset paper sizes (width x height in inches)
PAPER_SIZES = {
    "letter": (8.5, 11.0),
    "a4": (8.27, 11.69),
    "legal": (8.5, 14.0),
    "tabloid": (11.0, 17.0),
}

# Preset grid configs: (rows, cols, rotate_landscape)
GRID_PRESETS = {
    "mtg":       (4, 2, True),    # 8 cards/page, landscape rotation
    "mtg-3x3":   (3, 3, False),   # 9 cards/page, portrait
    "token":     (3, 3, False),   # 9 tokens/page
    "full-art":  (2, 2, False),   # 4 large cards/page
    "playtest":  (5, 3, False),   # 15 small cards/page
}


def rotate_image_landscape(img_path: Path, temp_dir: Path) -> Path:
    """Rotate image 90 degrees for landscape placement. Returns path to rotated image."""
    temp_path = temp_dir / f"rotated_{img_path.name}"
    if temp_path.exists():
        return temp_path

    with Image.open(img_path) as img:
        rotated = img.rotate(-90, expand=True)
        temp_path.parent.mkdir(parents=True, exist_ok=True)
        rotated.save(temp_path)

    return temp_path


def create_grid(
    image_paths: List[str],
    rows: int = 4,
    columns: int = 2,
    output_path: str = "outputs/grid.pptx",
    slide_width: float = 8.5,
    slide_height: float = 11.0,
    margin: float = 0.3,
    spacing: float = 0.15,
    rotate: bool = True,
    bg_color: str = "white",
) -> Dict[str, Any]:
    """
    Create presentation with images in a grid layout.

    Args:
        image_paths: List of image file paths
        rows: Rows per slide
        columns: Columns per slide
        output_path: Output PPTX path
        slide_width: Paper width in inches
        slide_height: Paper height in inches
        margin: Page margin in inches
        spacing: Gap between cells in inches
        rotate: Rotate images 90 degrees for landscape placement
        bg_color: Background color ("white" or "black")

    Returns:
        Dict with path, slide_count, images_per_slide, card dimensions
    """
    if rows <= 0 or columns <= 0:
        print(f"Error: rows and columns must be > 0")
        sys.exit(1)

    if not image_paths:
        print(f"Error: no image paths provided")
        sys.exit(1)

    # Filter to existing files
    valid_paths = []
    missing = []
    for p in image_paths:
        if Path(p).exists():
            valid_paths.append(Path(p))
        else:
            missing.append(p)

    if missing:
        print(f"Warning: {len(missing)} images not found, skipping")

    if not valid_paths:
        print("Error: no valid image files found")
        sys.exit(1)

    # Calculate cell dimensions
    available_width = slide_width - (2 * margin) - ((columns - 1) * spacing)
    available_height = slide_height - (2 * margin) - ((rows - 1) * spacing)
    cell_width = available_width / columns
    cell_height = available_height / rows

    images_per_slide = rows * columns
    slide_count = math.ceil(len(valid_paths) / images_per_slide)

    print(f"Grid: {columns}x{rows} ({images_per_slide} per page)")
    print(f"Paper: {slide_width}\" x {slide_height}\"")
    print(f"Cell size: {cell_width:.2f}\" x {cell_height:.2f}\"")
    print(f"Images: {len(valid_paths)} across {slide_count} pages")
    if rotate:
        print(f"Rotation: 90 degrees (landscape)")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    blank_layout = prs.slide_layouts[6]

    # Temp dir for rotated images
    temp_dir = Path("images/.temp_rotated")
    if rotate:
        temp_dir.mkdir(parents=True, exist_ok=True)

    bg_rgb = RGBColor(255, 255, 255) if bg_color == "white" else RGBColor(0, 0, 0)
    placed = 0

    for slide_num in range(slide_count):
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

        start_idx = slide_num * images_per_slide
        end_idx = min(start_idx + images_per_slide, len(valid_paths))

        for i, card_idx in enumerate(range(start_idx, end_idx)):
            img_path = valid_paths[card_idx]
            row = i // columns
            col = i % columns

            left = margin + col * (cell_width + spacing)
            top = margin + row * (cell_height + spacing)

            # Rotate if needed
            if rotate:
                use_path = rotate_image_landscape(img_path, temp_dir)
            else:
                use_path = img_path

            try:
                # Fit image preserving aspect ratio
                with Image.open(use_path) as img:
                    img_aspect = img.width / img.height
                cell_aspect = cell_width / cell_height

                if img_aspect > cell_aspect:
                    w = cell_width
                    h = cell_width / img_aspect
                else:
                    h = cell_height
                    w = cell_height * img_aspect

                # Center in cell
                x = left + (cell_width - w) / 2
                y = top + (cell_height - h) / 2

                slide.shapes.add_picture(
                    str(use_path),
                    Inches(x), Inches(y),
                    width=Inches(w), height=Inches(h)
                )
                placed += 1

                # Add crop marks at card corners (0.1" lines)
                mark_len = Inches(0.1)
                mark_width = Inches(0.004)  # ~0.5pt
                mark_color = RGBColor(0, 0, 0)
                corners = [
                    (x, y),                     # top-left
                    (x + w, y),                 # top-right
                    (x, y + h),                 # bottom-left
                    (x + w, y + h),             # bottom-right
                ]
                for cx, cy in corners:
                    # Horizontal mark
                    from pptx.util import Emu
                    from pptx.enum.shapes import MSO_SHAPE
                    hline = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cx - 0.05 if cx == x else cx),
                        Inches(cy - 0.002),
                        mark_len, mark_width,
                    )
                    hline.fill.solid()
                    hline.fill.fore_color.rgb = mark_color
                    hline.line.fill.background()
                    # Vertical mark
                    vline = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cx - 0.002),
                        Inches(cy - 0.05 if cy == y else cy),
                        mark_width, mark_len,
                    )
                    vline.fill.solid()
                    vline.fill.fore_color.rgb = mark_color
                    vline.line.fill.background()

            except Exception as e:
                print(f"  Warning: failed to place {img_path.name}: {e}")

    # Save
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))

    # Clean up temp rotated images
    if rotate and temp_dir.exists():
        for f in temp_dir.iterdir():
            f.unlink()
        temp_dir.rmdir()

    print(f"\nSaved: {output_file}")
    print(f"Placed: {placed}/{len(valid_paths)} images")

    return {
        "path": str(output_file),
        "slide_count": slide_count,
        "images_per_slide": images_per_slide,
        "images_placed": placed,
        "cell_width_inches": round(cell_width, 2),
        "cell_height_inches": round(cell_height, 2),
    }


def convert_to_pdf(pptx_path: str, output_dir: str = None) -> Optional[str]:
    """Convert PPTX to PDF via LibreOffice."""
    pptx = Path(pptx_path)
    out_dir = output_dir or str(pptx.parent)

    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", out_dir, str(pptx)],
            check=True, capture_output=True, timeout=60
        )
        pdf_path = Path(out_dir) / f"{pptx.stem}.pdf"
        if pdf_path.exists():
            print(f"PDF: {pdf_path}")
            return str(pdf_path)
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
        print("PDF conversion skipped (LibreOffice not available)")

    return None


def main():
    parser = argparse.ArgumentParser(
        description="Generate print-ready grid layouts from images",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Presets:
  mtg       2x4 landscape, 8 cards/page (default)
  mtg-3x3   3x3 portrait, 9 cards/page
  token     3x3 portrait, 9 tokens/page
  full-art  2x2 portrait, 4 large cards/page
  playtest  3x5 portrait, 15 small cards/page

Examples:
  %(prog)s images/*.jpg
  %(prog)s --preset mtg-3x3 images/*.jpg
  %(prog)s --rows 3 --cols 3 --no-rotate images/*.png
  %(prog)s --paper a4 --output outputs/my_grid.pptx images/*.jpg
        """
    )

    parser.add_argument("images", nargs="+", help="Image file paths")
    parser.add_argument("-o", "--output", default="outputs/grid.pptx",
                        help="Output PPTX path (default: outputs/grid.pptx)")
    parser.add_argument("--preset", choices=GRID_PRESETS.keys(),
                        help="Use a preset grid configuration")
    parser.add_argument("--rows", type=int, help="Rows per page")
    parser.add_argument("--cols", type=int, help="Columns per page")
    parser.add_argument("--paper", choices=PAPER_SIZES.keys(), default="letter",
                        help="Paper size (default: letter)")
    parser.add_argument("--width", type=float, help="Custom paper width in inches")
    parser.add_argument("--height", type=float, help="Custom paper height in inches")
    parser.add_argument("--margin", type=float, default=0.3,
                        help="Page margin in inches (default: 0.3)")
    parser.add_argument("--spacing", type=float, default=0.15,
                        help="Gap between cells in inches (default: 0.15)")
    parser.add_argument("--no-rotate", action="store_true",
                        help="Don't rotate images (keep portrait orientation)")
    parser.add_argument("--bg", choices=["white", "black"], default="white",
                        help="Background color (default: white)")
    parser.add_argument("--pdf", action="store_true",
                        help="Also convert to PDF via LibreOffice")
    parser.add_argument("--json", action="store_true",
                        help="Output result as JSON")

    args = parser.parse_args()

    # Resolve grid config
    if args.preset:
        rows, cols, rotate = GRID_PRESETS[args.preset]
    else:
        rows = args.rows or 4
        cols = args.cols or 2
        rotate = True

    # --no-rotate overrides preset
    if args.no_rotate:
        rotate = False

    # Resolve paper size
    if args.width and args.height:
        slide_width, slide_height = args.width, args.height
    else:
        slide_width, slide_height = PAPER_SIZES[args.paper]

    result = create_grid(
        image_paths=args.images,
        rows=rows,
        columns=cols,
        output_path=args.output,
        slide_width=slide_width,
        slide_height=slide_height,
        margin=args.margin,
        spacing=args.spacing,
        rotate=rotate,
        bg_color=args.bg,
    )

    if args.pdf:
        pdf_path = convert_to_pdf(args.output)
        if pdf_path:
            result["pdf_path"] = pdf_path

    if args.json:
        print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
