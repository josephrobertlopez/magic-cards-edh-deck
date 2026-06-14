#!/usr/bin/env python3
"""
Generate MTG deck on US Letter paper (8.5" × 11") with 8 cards per page.
Cards are placed in landscape orientation (3.5"w × 2.5"h) in 2×4 grid.
"""
import sys
import time
import requests
from pathlib import Path
import subprocess
import math
import json
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image

METADATA_CACHE_PATH = "images/card_metadata.json"

# Scryfall blocks the default python-requests User-Agent (400). Their API
# requires an explicit User-Agent + Accept header.
SCRYFALL_HEADERS = {
    "User-Agent": "magic-cards-edh-deck/1.0",
    "Accept": "application/json",
}

def safe_filename(name):
    """Sanitize card name for use as filename."""
    return name.replace(' ', '_').replace('/', '_').replace("'", "").replace(',', '')

def load_metadata_cache():
    """Load card metadata cache from JSON file."""
    if Path(METADATA_CACHE_PATH).exists():
        with open(METADATA_CACHE_PATH, 'r') as f:
            return json.load(f)
    return {}

def save_metadata_cache(cache):
    """Save card metadata cache to JSON file."""
    Path("images").mkdir(exist_ok=True)
    with open(METADATA_CACHE_PATH, 'w') as f:
        json.dump(cache, f, indent=2)

def fetch_card_from_scryfall(card_name):
    """Fetch card data from Scryfall API with AGGRESSIVE ERROR DETECTION."""
    # HOSTILE VALIDATION: Check for obvious comment pollution
    if card_name.startswith('#'):
        print(f"🚨 FATAL ERROR: Trying to fetch COMMENT as card: '{card_name}'")
        print(f"🚨 THIS SHOULD NEVER HAPPEN - PARSER FAILED!")
        raise ValueError(f"Comment pollution detected in API call: {card_name}")

    if not card_name.strip():
        print(f"🚨 FATAL ERROR: Empty card name passed to API")
        raise ValueError("Empty card name - parser validation failed")

    print(f"🌐 API REQUEST: Searching Scryfall for '{card_name}'")
    url = "https://api.scryfall.com/cards/named"
    params = {"fuzzy": card_name}

    backoff_delay = 2.0
    for attempt in range(5):
        try:
            response = requests.get(url, params=params, headers=SCRYFALL_HEADERS, timeout=10)

            # AGGRESSIVE STATUS LOGGING
            print(f"    📡 API Response: {response.status_code} for '{card_name}'")

            if response.status_code == 429:
                if attempt < 4:
                    print(f"    🔥 RATE LIMITED (attempt {attempt+1}/5): waiting {backoff_delay}s...")
                    time.sleep(backoff_delay)
                    backoff_delay = min(backoff_delay * 2, 30)
                    continue
                else:
                    print(f"  💀 TOTAL FAILURE: '{card_name}' - Rate limited after 5 retries")
                    print(f"  💀 WASTED {5} API CALLS ON: '{card_name}'")
                    return None

            if response.status_code == 404:
                print(f"  💀 CARD NOT FOUND: '{card_name}' (404 - likely typo or comment)")
                return None

            response.raise_for_status()
            card_data = response.json()
            print(f"  ✅ SUCCESS: Found '{card_data.get('name', card_name)}'")
            return card_data

        except Exception as e:
            print(f"  🚨 EXCEPTION on '{card_name}': {type(e).__name__}: {e}")
            if "429" not in str(e):
                print(f"  💀 COMPLETE FAILURE: '{card_name}' - {e}")
            return None
    return None

def download_image(url, output_path):
    """Download image from URL."""
    try:
        response = requests.get(url, headers=SCRYFALL_HEADERS, timeout=30)
        response.raise_for_status()

        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'wb') as f:
            f.write(response.content)
        return True
    except Exception as e:
        print(f"  ❌ Failed to download image: {e}")
        return False

def create_double_sided_composite(front_path, back_path, output_path):
    """
    Create a side-by-side composite of double-sided card for printing.
    Both faces on one card - turn sideways and stack them.
    """
    try:
        with Image.open(front_path) as front_img, Image.open(back_path) as back_img:
            # Standard MTG card dimensions
            card_width = front_img.width
            card_height = front_img.height

            # Create composite: side-by-side layout (double width)
            # Left = Front face, Right = Back face (flipped for proper stacking)
            composite_width = card_width * 2
            composite_height = card_height

            composite = Image.new('RGB', (composite_width, composite_height), 'white')

            # Place front face on left
            composite.paste(front_img, (0, 0))

            # Place back face on right (flip horizontally for proper card stacking)
            back_flipped = back_img.transpose(Image.FLIP_LEFT_RIGHT)
            composite.paste(back_flipped, (card_width, 0))

            # Save composite
            output_path.parent.mkdir(parents=True, exist_ok=True)
            composite.save(output_path, 'JPEG', quality=95)
            return True

    except Exception as e:
        print(f"  💀 COMPOSITE FAILED: {e}")
        return False

def get_letter_size_card_positions():
    """
    Calculate 8 card positions for US Letter paper (8.5" × 11").
    Returns 2×4 grid with cards in landscape orientation, centered on page.
    """
    # Letter paper dimensions
    page_width = 8.5
    page_height = 11.0

    # Card dimensions in landscape (rotated 90°)
    # 2% larger than sleeve-compatible size for better visibility when printing
    # (2.60" × 3.62" = 2% larger than 2.55" × 3.55")
    card_width = 3.62  # MTG card height becomes width
    card_height = 2.60  # MTG card width becomes height

    # Gap between cards
    gap_horizontal = 0.2
    gap_vertical = 0.2

    # Calculate total dimensions needed
    total_width = 2 * card_width + gap_horizontal
    total_height = 4 * card_height + 3 * gap_vertical

    # Center on page
    margin_left = (page_width - total_width) / 2
    margin_top = (page_height - total_height) / 2

    positions = []

    # 2 columns × 4 rows
    for row in range(4):
        for col in range(2):
            left = margin_left + col * (card_width + gap_horizontal)
            top = margin_top + row * (card_height + gap_vertical)

            positions.append({
                'left': left,
                'top': top,
                'width': card_width,
                'height': card_height
            })

    return positions

def place_card_landscape(slide, card_path, slot_info):
    """
    Place card image in landscape orientation (rotated 90°).
    """
    try:
        if not card_path.exists():
            print(f"    ⚠️  Image not found: {card_path}")
            return False

        with Image.open(card_path) as img:
            # Rotate image 90° for landscape placement
            img_rotated = img.rotate(-90, expand=True)

            # Calculate aspect ratios
            slot_aspect = slot_info['width'] / slot_info['height']
            img_aspect = img_rotated.width / img_rotated.height

            # Preserve aspect ratio while fitting to slot
            if img_aspect > slot_aspect:
                new_width = slot_info['width']
                new_height = slot_info['width'] / img_aspect
            else:
                new_height = slot_info['height']
                new_width = slot_info['height'] * img_aspect

            # Center in slot
            left = slot_info['left'] + (slot_info['width'] - new_width) / 2
            top = slot_info['top'] + (slot_info['height'] - new_height) / 2

            # Save rotated image temporarily
            temp_path = card_path.parent / f"temp_rotated_{card_path.name}"
            img_rotated.save(temp_path)

        # Add picture to slide
        slide.shapes.add_picture(
            str(temp_path),
            Inches(left),
            Inches(top),
            width=Inches(new_width),
            height=Inches(new_height)
        )

        # Clean up temp file
        temp_path.unlink()

        return True

    except Exception as e:
        print(f"    ❌ Failed to place card: {e}")
        return False

def create_letter_size_deck(card_images, output_file):
    """
    Create deck presentation on US Letter paper (8.5" × 11").
    """
    print(f"📄 Creating Letter size presentation (8.5\" × 11\")")

    # Get card positions for Letter paper
    card_positions = get_letter_size_card_positions()
    cards_per_page = len(card_positions)  # Should be 8
    slides_needed = math.ceil(len(card_images) / cards_per_page)

    print(f"📐 Layout: 2 columns × 4 rows (8 cards per page)")
    print(f"📏 Cards in landscape: 3.5\" × 2.5\" each")
    print(f"📄 Creating {slides_needed} slides")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11.0)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Process cards in groups of 8
    cards_placed = 0
    for slide_num in range(slides_needed):
        # Add new slide
        slide = prs.slides.add_slide(blank_layout)

        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Calculate card indices for this slide
        start_idx = slide_num * cards_per_page
        end_idx = min(start_idx + cards_per_page, len(card_images))

        print(f"\n📄 Slide {slide_num + 1}: Cards {start_idx + 1} to {end_idx}")

        # Place cards using calculated positions
        for i, card_idx in enumerate(range(start_idx, end_idx)):
            if i >= len(card_positions):
                break

            card_path = card_images[card_idx]
            slot = card_positions[i]

            card_name = card_path.stem
            print(f"  [{i+1}] {card_name}")
            if place_card_landscape(slide, card_path, slot):
                print(f"      ✅ Placed at ({slot['left']:.2f}\", {slot['top']:.2f}\") - {slot['width']:.2f}\" × {slot['height']:.2f}\"")
                cards_placed += 1
            else:
                print(f"      ❌ Failed to place")

    # Save presentation
    prs.save(output_file)
    print(f"\n💾 Saved presentation: {output_file}")
    print(f"📊 Stats: {cards_placed} cards across {slides_needed} slides")
    print(f"📏 Paper size: 8.5\" × 11\" (US Letter)")
    print(f"🎴 Card layout: 2×4 grid, landscape orientation")

    return True

def main(card_list_file, output_dir="outputs"):
    """Generate deck on Letter size paper."""

    print("🎴 MTG Deck Generator (US Letter Size - 8 cards per page)")
    print("=" * 70)

    # Read card list
    card_list_path = Path(card_list_file)
    if not card_list_path.exists():
        print(f"❌ Card list not found: {card_list_file}")
        return 1

    # ADVERSARIAL PARSING WITH AGGRESSIVE VALIDATION
    print("🔍 PARSING DECKLIST WITH HOSTILE VALIDATION...")

    with open(card_list_path, 'r') as f:
        all_lines = [line.strip() for line in f if line.strip()]

    # SCREAM ABOUT COMMENTS vs CARDS
    comments_found = []
    card_names = []

    for line_num, line in enumerate(all_lines, 1):
        if line.startswith('#'):
            comments_found.append(line)
            print(f"🚨 LINE {line_num}: COMMENT DETECTED (IGNORING): '{line}'")
        elif line:
            card_names.append(line)
            print(f"✅ LINE {line_num}: CARD DETECTED: '{line}'")

    # AGGRESSIVE LOGGING OF PARSING RESULTS
    print("=" * 70)
    print(f"🔥 PARSING RESULTS (ADVERSARIAL VALIDATION):")
    print(f"📊 TOTAL LINES PROCESSED: {len(all_lines)}")
    print(f"💀 COMMENTS FOUND: {len(comments_found)}")
    print(f"🎯 ACTUAL CARDS FOUND: {len(card_names)}")

    if comments_found:
        print("🚨 COMMENTS THAT WOULD HAVE BROKEN OLD PARSER:")
        for comment in comments_found[:5]:  # Show first 5
            print(f"   💀 '{comment}'")
        if len(comments_found) > 5:
            print(f"   💀 ... and {len(comments_found) - 5} more comments")

    # SCREAM IF EMPTY DECK
    if not card_names:
        print("🚨 FATAL: NO CARDS FOUND IN DECKLIST!")
        print("🚨 Check your file format - only comments detected!")
        return 1

    print("=" * 70)
    print(f"📋 VALIDATED: {len(card_names)} actual cards ready for processing")
    print()

    # Create output directories
    images_dir = Path("images")
    images_dir.mkdir(exist_ok=True)
    output_path = Path(output_dir)
    output_path.mkdir(exist_ok=True)

    # Load metadata cache
    metadata_cache = load_metadata_cache()

    # ADVERSARIAL API PULL TRACKING
    print("🔍 Fetching card data from Scryfall...")
    successful_images = []
    api_failures = []
    cache_hits = []
    download_failures = []
    partial_failures = []

    for i, card_name in enumerate(card_names, 1):
        print(f"[{i}/{len(card_names)}] {card_name}")

        # Check if metadata exists and all images are cached
        if card_name in metadata_cache:
            card_metadata = metadata_cache[card_name]
            faces = card_metadata.get('faces', [])

            # Check if all face images exist on disk
            all_cached = True
            for face_name in faces:
                safe_name = safe_filename(face_name)
                image_path = images_dir / f"{safe_name}.jpg"
                if not image_path.exists():
                    all_cached = False
                    break

            # If all images cached, add each face as its own upright card (no mirrored composite)
            if all_cached:
                for face_name in faces:
                    safe_name = safe_filename(face_name)
                    image_path = images_dir / f"{safe_name}.jpg"
                    print(f"  ✓ {face_name} (cached)")
                    successful_images.append(image_path)
                cache_hits.append(card_name)
                continue  # Skip API call and delay

        # Fallback: image exists on disk but no metadata (e.g. downloaded externally)
        fallback_name = safe_filename(card_name)
        fallback_path = images_dir / f"{fallback_name}.jpg"
        if fallback_path.exists():
            metadata_cache[card_name] = {
                'type': 'single',
                'faces': [card_name],
                'fetched': time.strftime("%Y-%m-%d")
            }
            save_metadata_cache(metadata_cache)
            print(f"  ✓ {card_name} (cached)")
            successful_images.append(fallback_path)
            cache_hits.append(card_name)
            continue

        # Fetch card data from Scryfall (needed if metadata AND image both missing)
        card_data = fetch_card_from_scryfall(card_name)
        if not card_data:
            api_failures.append(card_name)
            continue

        # Store metadata in cache
        if 'card_faces' in card_data:
            faces_list = [face.get('name', card_name) for face in card_data['card_faces']]
        else:
            faces_list = [card_name]

        metadata_cache[card_name] = {
            'type': 'double_faced' if 'card_faces' in card_data else 'single',
            'faces': faces_list,
            'fetched': time.strftime("%Y-%m-%d")
        }
        save_metadata_cache(metadata_cache)

        # Double-faced cards (MDFC / transform): each face with its own art prints as a separate upright card
        printable_faces = [f for f in card_data.get('card_faces', []) if f.get('image_uris', {}).get('normal')]
        if len(printable_faces) >= 2:
            print(f"  🔄 DOUBLE-FACED CARD: {len(printable_faces)} faces -> separate proxy cards")

            for face_idx, face in enumerate(printable_faces):
                face_name = face.get('name', card_name)
                face_type = "Front" if face_idx == 0 else "Back"
                image_url = face['image_uris']['normal']

                safe_name = safe_filename(face_name)
                image_path = images_dir / f"{safe_name}.jpg"

                if image_path.exists():
                    print(f"  ✓ {face_type}: {face_name} (cached)")
                    successful_images.append(image_path)
                else:
                    if download_image(image_url, image_path):
                        print(f"  ✓ {face_type}: {face_name}")
                        successful_images.append(image_path)
                    else:
                        download_failures.append(f"{card_name} ({face_type})")
                        print(f"  💀 DOWNLOAD FAILED: {face_type}: {face_name}")
                time.sleep(0.1)  # Rate limiting between faces

        elif 'image_uris' in card_data:
            # Single-faced card
            image_url = card_data['image_uris'].get('normal')

            if not image_url:
                print(f"  ⚠️  No image available")
                partial_failures.append(card_name)
                continue

            # Download image
            safe_name = safe_filename(card_name)
            image_path = images_dir / f"{safe_name}.jpg"

            if image_path.exists():
                print(f"  ✓ Already downloaded")
                successful_images.append(image_path)
            else:
                if download_image(image_url, image_path):
                    print(f"  ✓ Downloaded")
                    successful_images.append(image_path)
                else:
                    download_failures.append(card_name)
                    print(f"  💀 DOWNLOAD FAILED: {card_name}")
        else:
            print(f"  ⚠️  No image available")
            partial_failures.append(card_name)
            continue

        # Rate limiting (0.15s for API calls, 0s for cached cards already skipped)
        time.sleep(0.15)

    print()
    print("=" * 70)
    print("🚨 ADVERSARIAL API PULL ANALYSIS (FAILURE DETECTION)")
    print("=" * 70)

    # Calculate success metrics
    total_requested = len(card_names)
    total_successful = len(successful_images)
    success_rate = (total_successful / total_requested * 100) if total_requested > 0 else 0

    print(f"📊 TOTAL CARDS REQUESTED: {total_requested}")
    print(f"✅ SUCCESSFUL IMAGES: {total_successful}")
    print(f"📈 SUCCESS RATE: {success_rate:.1f}%")
    print()

    # SCREAM ABOUT CACHE EFFICIENCY
    cache_rate = (len(cache_hits) / total_requested * 100) if total_requested > 0 else 0
    if cache_rate > 80:
        print(f"🎯 CACHE HIT RATE: {cache_rate:.1f}% (EXCELLENT)")
    elif cache_rate > 50:
        print(f"⚠️  CACHE HIT RATE: {cache_rate:.1f}% (MODERATE)")
    else:
        print(f"🚨 CACHE HIT RATE: {cache_rate:.1f}% (POOR - HEAVY API USAGE)")

    # SCREAM ABOUT FAILURES
    if api_failures:
        print(f"🚨 API FAILURES ({len(api_failures)}):")
        for failure in api_failures[:5]:
            print(f"   💀 '{failure}' - API call failed")
        if len(api_failures) > 5:
            print(f"   💀 ... and {len(api_failures) - 5} more API failures")

    if download_failures:
        print(f"🚨 DOWNLOAD FAILURES ({len(download_failures)}):")
        for failure in download_failures[:5]:
            print(f"   💀 '{failure}' - Image download failed")
        if len(download_failures) > 5:
            print(f"   💀 ... and {len(download_failures) - 5} more download failures")

    if partial_failures:
        print(f"🚨 PARTIAL FAILURES ({len(partial_failures)}):")
        for failure in partial_failures[:5]:
            print(f"   ⚠️  '{failure}' - No image available from Scryfall")
        if len(partial_failures) > 5:
            print(f"   ⚠️  ... and {len(partial_failures) - 5} more partial failures")

    # OVERALL STATUS
    if success_rate >= 95:
        print("✅ OVERALL STATUS: EXCELLENT")
    elif success_rate >= 85:
        print("⚠️  OVERALL STATUS: ACCEPTABLE")
    else:
        print("🚨 OVERALL STATUS: PROBLEMATIC - INVESTIGATE FAILURES")

    print("=" * 70)

    if len(successful_images) == 0:
        print("❌ No cards downloaded, cannot generate deck")
        return 1

    # Generate PPTX
    print("📄 Generating Letter size deck...")
    pptx_path = output_path / "deck_letter_size.pptx"

    if not create_letter_size_deck(successful_images, pptx_path):
        return 1

    print()

    # Convert to PDF
    print("📄 Converting to PDF...")
    pdf_path = output_path / "deck_letter_size.pdf"

    try:
        subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_path),
            str(pptx_path)
        ], check=True, capture_output=True, timeout=60)
        print(f"  ✓ {pdf_path}")
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as e:
        print(f"  ⚠️  PDF conversion failed (LibreOffice not available)")
        print(f"  ℹ️  You can convert {pptx_path} to PDF manually")

    print()
    print("=" * 70)
    print("🎉 Deck generation complete!")
    print(f"📁 PPTX: {pptx_path}")
    if pdf_path.exists():
        print(f"📁 PDF:  {pdf_path}")
    print(f"🎴 {len(successful_images)} cards on US Letter paper (8.5\" × 11\")")
    print(f"📏 8 cards per page in 2×4 grid (landscape)")
    print(f"✂️  Ready to print and cut!")

    return 0

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 generate_deck_letter_size.py <card_list_file>")
        print()
        print("Example:")
        print("  python3 generate_deck_letter_size.py decklists/user_deck.txt")
        sys.exit(1)

    sys.exit(main(sys.argv[1]))
