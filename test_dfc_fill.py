#!/usr/bin/env python3
import os
import json
import math
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image
from pathlib import Path

def get_template_slot_positions(template_slide):
    """
    Extract card slot positions from template (FR-006)
    Returns list of slot dictionaries with left, top, width, height
    """
    card_slots = []

    for shape in template_slide.shapes:
        # Look for rectangular shapes that are card-sized
        if (shape.shape_type in [1, 13] and  # Rectangle or Picture
            shape.width.inches > 1.0 and
            shape.height.inches > 1.0):

            card_slots.append({
                'left': shape.left.inches,
                'top': shape.top.inches,
                'width': shape.width.inches,
                'height': shape.height.inches
            })

    # Sort slots by position (top to bottom, left to right)
    card_slots.sort(key=lambda x: (x['top'], x['left']))

    return card_slots

def place_card_in_slot(slide, card_path, slot_info):
    """
    Place card image in slot with aspect ratio preservation (FR-005, FR-007)
    Handles vertical/horizontal orientation based on slot dimensions
    """
    try:
        # Validate card image exists
        if not os.path.exists(card_path):
            print(f"    ⚠️  Image not found: {card_path}")
            return False

        with Image.open(card_path) as img:
            # Calculate aspect ratios
            slot_aspect = slot_info['width'] / slot_info['height']
            img_aspect = img.width / img.height

            # Determine if slot is horizontal (FR-007)
            is_horizontal_slot = slot_aspect > 1.0

            # Preserve aspect ratio while fitting to slot (FR-005)
            if img_aspect > slot_aspect:
                # Image is wider, fit to width
                new_width = slot_info['width']
                new_height = slot_info['width'] / img_aspect
            else:
                # Image is taller, fit to height
                new_height = slot_info['height']
                new_width = slot_info['height'] * img_aspect

            # Center the image in the slot
            left = slot_info['left'] + (slot_info['width'] - new_width) / 2
            top = slot_info['top'] + (slot_info['height'] - new_height) / 2

            # Resize for quality (150 DPI equivalent)
            target_width = int(new_width * 150)
            target_height = int(new_height * 150)

            resized_img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)

            # Rotate if horizontal slot and card is vertical
            if is_horizontal_slot and img_aspect < 1.0:
                resized_img = resized_img.rotate(90, expand=True)

            # Save temporary file
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as temp_file:
                temp_path = temp_file.name
                resized_img.save(temp_path, "JPEG", quality=90)

        # Add picture to slide
        slide.shapes.add_picture(
            temp_path,
            Inches(left),
            Inches(top),
            Inches(new_width),
            Inches(new_height)
        )

        # Clean up temp file
        os.remove(temp_path)

        return True

    except Exception as e:
        print(f"    ❌ Failed to place card: {e}")
        return False

def create_presentation_from_manifest(template_file, output_file, manifest_path):
    """
    Create presentation using template pattern with DFC support
    Handles both single-path and multi-path (DFC) cards
    """

    # Validate template exists
    if not os.path.exists(template_file):
        print(f"❌ Error: Template file not found: {template_file}")
        return False

    # Load manifest
    if not os.path.exists(manifest_path):
        print(f"❌ Error: Manifest not found: {manifest_path}")
        return False

    with open(manifest_path, 'r') as f:
        manifest = json.load(f)

    # Flatten cards: Handle DFC cards with multiple faces
    card_images = []
    for card in manifest['cards']:
        if card['status'] != 'success':
            continue

        # Check if this is a multi-face card
        if 'paths' in card and card.get('faces', 0) > 1:
            # Add each face as a separate image
            for i, face_path in enumerate(card['paths']):
                card_images.append({
                    'name': f"{card['name']} (Face {i+1})",
                    'path': face_path
                })
        else:
            # Single face card
            card_images.append({
                'name': card['name'],
                'path': card['path']
            })

    total_cards = len(card_images)
    print(f"🎴 Found {total_cards} card images to place (including DFC faces)")

    if total_cards == 0:
        print("❌ No successful cards in manifest")
        return False

    # Load template
    print(f"📋 Loading template: {template_file}")
    template_prs = Presentation(template_file)

    if len(template_prs.slides) == 0:
        print("❌ Template has no slides")
        return False

    # Get slot positions from first slide
    template_slide = template_prs.slides[0]
    slot_positions = get_template_slot_positions(template_slide)

    print(f"📐 Found {len(slot_positions)} card slots in template")

    if len(slot_positions) == 0:
        print("❌ No card slots found in template")
        return False

    cards_per_slide = len(slot_positions)
    slides_needed = math.ceil(total_cards / cards_per_slide)

    print(f"📄 Creating {slides_needed} slides ({cards_per_slide} cards per slide)")

    # Create new presentation from template to preserve dimensions
    prs = Presentation(template_file)

    # Remove template slides, keep only the slide master/layouts
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    blank_layout = prs.slide_layouts[6]  # Blank layout from template

    # Process cards in groups
    for slide_num in range(slides_needed):
        # Add new slide
        slide = prs.slides.add_slide(blank_layout)

        # Black background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Calculate card indices for this slide
        start_idx = slide_num * cards_per_slide
        end_idx = min(start_idx + cards_per_slide, total_cards)

        print(f"\n📄 Slide {slide_num + 1}: Cards {start_idx + 1} to {end_idx}")

        # Place cards using template positions
        for i, card_idx in enumerate(range(start_idx, end_idx)):
            if i >= len(slot_positions):
                break

            card_info = card_images[card_idx]
            card_path = card_info['path']
            card_name = card_info['name']
            slot = slot_positions[i]

            print(f"  [{i+1}] {card_name}")
            if place_card_in_slot(slide, card_path, slot):
                print(f"      ✅ Placed in slot")
            else:
                print(f"      ❌ Failed to place")

        # Leave empty slots blank if cards exhausted
        remaining_slots = len(slot_positions) - (end_idx - start_idx)
        if remaining_slots > 0:
            print(f"  ℹ️  {remaining_slots} empty slots left blank")

    # Save presentation
    prs.save(output_file)
    print(f"\n💾 Saved presentation: {output_file}")
    print(f"📊 Stats: {total_cards} card images across {slides_needed} slides")

    return True

if __name__ == "__main__":
    print(f"🎴 Testing DFC Card Placement")
    print("=" * 70)

    template_file = "decks.pptx"
    output_file = "tests/dfc_output.pptx"
    manifest_path = ".claude/state/dfc_test_manifest.json"

    success = create_presentation_from_manifest(template_file, output_file, manifest_path)

    if success:
        print("\n" + "=" * 70)
        print(f"🎉 DFC presentation created successfully!")
        print(f"📁 File: {output_file}")
    else:
        print("\n❌ Failed to create DFC presentation")
