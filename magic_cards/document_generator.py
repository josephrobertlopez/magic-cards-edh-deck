"""Document generator module for creating PPTX presentations."""

import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from pathlib import Path


def find_card_slots(prs):
    """Find all card placeholder slots in presentation."""
    slots = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Look for shapes that could hold cards
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or 
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE or
                shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER):
                
                # Check if it's roughly card-sized
                width_inches = shape.width / 914400  # Convert EMU to inches
                height_inches = shape.height / 914400
                
                # Accept shapes between 1" and 5" in both dimensions
                if (1.0 <= width_inches <= 5.0 and 1.0 <= height_inches <= 5.0):
                    slots.append({
                        'slide': slide,
                        'slide_idx': slide_idx,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'width_inches': width_inches,
                        'height_inches': height_inches
                    })
    
    return slots


def resize_image_for_slot(image_path, target_width_inches, target_height_inches):
    """Resize image to fit card slot."""
    temp_dir = Path("temp_resized")
    temp_dir.mkdir(exist_ok=True)
    
    with Image.open(image_path) as img:
        # Calculate target size in pixels (96 DPI)
        target_width_px = int(target_width_inches * 96)
        target_height_px = int(target_height_inches * 96)
        
        # Resize maintaining aspect ratio
        img.thumbnail((target_width_px, target_height_px), Image.Resampling.LANCZOS)
        
        # Save resized image
        filename = Path(image_path).name
        temp_path = temp_dir / f"resized_{filename}"
        img.save(temp_path, "PNG")
        
        return str(temp_path)


def generate_document(manifest_path, template_path, output_path):
    """
    Generate PowerPoint presentation from manifest and template.
    
    Args:
        manifest_path: Path to manifest JSON with card image info
        template_path: Path to PowerPoint template file
        output_path: Path for output PPTX file
        
    Returns:
        Output path if successful, None otherwise
    """
    print(f"🎨 Generating presentation...")
    print(f"  📄 Template: {template_path}")
    print(f"  📋 Manifest: {manifest_path}")
    
    # Load manifest
    with open(manifest_path, 'r') as f:
        manifest = json.load(f)
    
    # Get successful card images
    card_images = [
        item['filepath'] 
        for item in manifest.get('items', [])
        if item.get('status') == 'success' and 'filepath' in item
    ]
    
    if not card_images:
        print("  ❌ No card images found in manifest")
        return None
    
    print(f"  🃏 Found {len(card_images)} card images")
    
    # Load presentation template
    prs = Presentation(template_path)
    print(f"  ✅ Loaded template with {len(prs.slides)} slides")
    
    # Find card slots
    slots = find_card_slots(prs)
    print(f"  🎴 Found {len(slots)} card slots")
    
    # Place cards in slots
    placed_cards = 0
    for idx, slot in enumerate(slots):
        if idx >= len(card_images):
            print(f"  ⚠️  Ran out of card images at slot {idx + 1}")
            break
        
        card_image = card_images[idx]
        
        try:
            # Resize image to fit slot
            resized_image = resize_image_for_slot(
                card_image,
                slot['width_inches'],
                slot['height_inches']
            )
            
            # Add image to slide at slot position
            slide = slot['slide']
            picture = slide.shapes.add_picture(
                resized_image,
                slot['left'],
                slot['top'],
                slot['width'],
                slot['height']
            )
            
            placed_cards += 1
            
        except Exception as e:
            print(f"  ❌ Failed to place card {idx + 1}: {e}")
            continue
    
    # Ensure output directory exists
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    
    # Save presentation
    prs.save(output_path)
    
    print(f"  ✅ Placed {placed_cards} cards")
    print(f"  💾 Saved: {output_path}")
    
    # Clean up temp files
    import shutil
    temp_dir = Path("temp_resized")
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
    
    return output_path
